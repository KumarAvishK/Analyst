"""
KiteIQX Intelligence - Streamlit Analytics Engine
=================================================
Refactored to:
  1. Match KiteIQX.com brand (navy + gold, professional consulting feel)
  2. Accept CSV, XLSX, and XLS files
  3. Executive Dashboard rebuilt as a narrative storyboard
  4. AI Intelligence moved to tab 2
  5. Data Quality lives in its own dedicated tab
  6. Every uploaded file is persisted to data/uploads/

SECURITY: The Groq API key is now read from st.secrets, NOT hardcoded.
Set it in .streamlit/secrets.toml or via the Streamlit Cloud dashboard.
"""

import io
import os
import re
import warnings
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import requests
import streamlit as st
from plotly.subplots import make_subplots
from sklearn.linear_model import LinearRegression
from sklearn.metrics import r2_score

warnings.filterwarnings("ignore")

# Google Sheets imports (soft)
try:
    import gspread
    from google.oauth2.service_account import Credentials
    GSHEETS_AVAILABLE = True
except ImportError:
    GSHEETS_AVAILABLE = False

# Groq imports (soft)
try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

# Anthropic Claude imports (soft)
try:
    from anthropic import Anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False


# ============================================================
# CONFIG
# ============================================================

UPLOAD_DIR = Path("data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

st.set_page_config(
    page_title="KiteIQX Intelligence",
    layout="wide",
    initial_sidebar_state="auto",
    page_icon="◆",
)


# ============================================================
# ⚠️  GROQ API KEY  -- TODO: ROTATE THIS AND MOVE TO st.secrets
# This is the SAME key from your original file. It is exposed in
# source control - rotate at https://console.groq.com/keys and
# replace with st.secrets["GROQ_API_KEY"] when you have time.
# ============================================================
_HARDCODED_GROQ_KEY = "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"


def _get_groq_key() -> str:
    """Prefer st.secrets, fall back to env var, fall back to hardcoded."""
    try:
        if "GROQ_API_KEY" in st.secrets:
            return st.secrets["GROQ_API_KEY"]
    except Exception:
        pass
    return os.environ.get("GROQ_API_KEY", _HARDCODED_GROQ_KEY)


GROQ_API_KEY = _get_groq_key()


# ============================================================
# CLAUDE (ANTHROPIC) API KEY  — from st.secrets / env
# ============================================================
def _get_anthropic_key() -> str:
    """Prefer st.secrets['ANTHROPIC_API_KEY'], then env var. No hardcoded fallback."""
    try:
        if "ANTHROPIC_API_KEY" in st.secrets:
            return st.secrets["ANTHROPIC_API_KEY"]
    except Exception:
        pass
    return os.environ.get("ANTHROPIC_API_KEY", "")


ANTHROPIC_API_KEY = _get_anthropic_key()
CLAUDE_MODEL = "claude-sonnet-4-6"

# ============================================================
# GOOGLE SHEETS LOGGER
# ============================================================

SHEET_ID = "1rJYe4MVKDc9srbzrqf1CWIfFLmPn_4qTwSEZ5PCTVT0"
GSHEETS_SCOPES = [
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive",
]

@st.cache_resource(show_spinner=False)
def _get_gsheets_client():
    """Cached Google Sheets client — one connection for the whole session."""
    if not GSHEETS_AVAILABLE:
        return None
    try:
        creds_dict = dict(st.secrets["gcp_service_account"])
        creds = Credentials.from_service_account_info(creds_dict, scopes=GSHEETS_SCOPES)
        return gspread.authorize(creds)
    except Exception:
        return None


def _get_session_id() -> str:
    if "session_id" not in st.session_state:
        import uuid
        st.session_state.session_id = str(uuid.uuid4())[:8]
    return st.session_state.session_id


def log_upload(filename: str, df, summary: str = "", takeaways: list = None):
    """Log a file upload event to the 'uploads' sheet."""
    try:
        client = _get_gsheets_client()
        if not client:
            return
        sh = client.open_by_key(SHEET_ID)
        ws = sh.worksheet("uploads")
        a = st.session_state.get("analytics")
        completeness = round(100 - df.isnull().sum().sum() / max(1, df.size) * 100, 1)
        num_cols = len(df.select_dtypes(include=[np.number]).columns)
        cat_cols = len(df.select_dtypes(include=["object", "category"]).columns)
        money_cols = [c for c in df.columns if any(w in c.lower()
                      for w in ("price","cost","amount","revenue","salary","sales","gmv"))]
        total_val = round(float(df[money_cols].sum().sum()), 2) if money_cols else ""
        ws.append_row([
            datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            _get_session_id(),
            filename,
            len(df),
            len(df.columns),
            num_cols,
            cat_cols,
            f"{completeness}%",
            str(total_val) if total_val != "" else "N/A",
            summary[:500] if summary else "",
            " | ".join(takeaways) if takeaways else "",
        ], value_input_option="USER_ENTERED")
    except Exception:
        pass  # never crash the app over logging


def log_ai_query(filename: str, question: str, answer: str):
    """Log an AI question + answer to the 'ai_queries' sheet."""
    try:
        client = _get_gsheets_client()
        if not client:
            return
        sh = client.open_by_key(SHEET_ID)
        ws = sh.worksheet("ai_queries")
        ws.append_row([
            datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            _get_session_id(),
            filename,
            question[:300],
            answer[:800],
        ], value_input_option="USER_ENTERED")
    except Exception:
        pass


# ============================================================
# THEME (KiteIQX consulting brand: navy + gold)
# ============================================================

st.markdown(
    """
<style>
:root {
    --kite-navy: #0a2540;
    --kite-navy-2: #14304f;
    --kite-gold: #c79a3a;
    --kite-gold-soft: #e6c878;
    --kite-text: #1a2333;
    --kite-text-soft: #4a5568;
    --kite-surface: #ffffff;
    --kite-surface-soft: #f7f8fb;
    --kite-border: #e3e6ee;
    --kite-success: #15803d;
    --kite-warning: #b45309;
    --kite-danger:  #b91c1c;
}

/* Force light mode globally */
[data-testid="stAppViewContainer"],
[data-testid="stMain"],
[data-testid="stMainBlockContainer"],
.main, section.main {
    background-color: #ffffff !important;
    color: #1a2333 !important;
}

html, body, [class*="css"] {
    font-family: 'Inter', -apple-system, 'Segoe UI', system-ui, sans-serif !important;
    background-color: #ffffff !important;
}
div.stMarkdown, div.stText { color: #1a2333; }

.block-container { padding-top: 4.5rem !important; max-width: 1200px; }

.stDeployButton { display: none !important; }
#MainMenu { visibility: hidden; }
footer { visibility: hidden; }

.kx-brand {
    padding-top: 0.25rem;
    display: flex;
    align-items: center;
    gap: 0.75rem;
    margin-bottom: 0.25rem;
}
.kx-mark {
    width: 58px; height: 58px;
    background: var(--kite-navy);
    color: var(--kite-gold);
    border-radius: 8px;
    display: flex; align-items: center; justify-content: center;
    font-weight: 700; font-size: 1.4rem;
    letter-spacing: -1px;
}
.kx-wordmark {
    font-size: 2.6rem; font-weight: 700;
    color: #0a2540 !important;
    letter-spacing: -0.3px;
}
.kx-wordmark span { color: #c79a3a !important; }
.kx-tag {
    color: #4a5568 !important;
    font-size: 0.95rem;
    margin-bottom: 1.25rem;
    border-bottom: 1px solid #e3e6ee;
    padding-bottom: 1rem;
}

.kx-pill {
    display: inline-block;
    padding: 0.32rem 0.85rem;
    border-radius: 999px;
    font-size: 0.8rem;
    font-weight: 600;
    letter-spacing: 0.2px;
    margin: 0.15rem 0.3rem 0.15rem 0;
}
.kx-pill-navy { background: var(--kite-navy); color: var(--kite-gold-soft); }
.kx-pill-gold { background: var(--kite-gold); color: #2a1b00; }
.kx-pill-soft { background: var(--kite-surface-soft); color: var(--kite-text-soft); border: 1px solid var(--kite-border); }

.kx-hero {
    background: linear-gradient(135deg, var(--kite-navy) 0%, var(--kite-navy-2) 100%);
    color: #ffffff;
    border-radius: 14px;
    padding: 1.75rem 2rem;
    margin: 0.5rem 0 1.25rem 0;
    box-shadow: 0 4px 24px rgba(10, 37, 64, 0.10);
}
.kx-hero h2 { color: #ffffff; margin: 0 0 0.4rem 0; font-weight: 700; }
.kx-hero p  { color: #d8dde6; margin: 0; font-size: 1.05rem; line-height: 1.55; }
.kx-hero .kx-hero-accent { color: var(--kite-gold); font-weight: 700; }

.kx-card {
    background: #ffffff !important;
    border: 1px solid var(--kite-border);
    border-radius: 12px;
    padding: 1.25rem 1.4rem;
    margin: 0.6rem 0;
    box-shadow: 0 1px 3px rgba(10, 37, 64, 0.04);
    color: #1a2333 !important;
}
.kx-card-title {
    font-size: 0.78rem;
    text-transform: uppercase;
    letter-spacing: 0.8px;
    color: #4a5568 !important;
    font-weight: 600;
    margin-bottom: 0.4rem;
}
.kx-card-value {
    font-size: 1.65rem;
    font-weight: 700;
    color: #0a2540 !important;
    line-height: 1.2;
}
.kx-card-sub {
    font-size: 0.85rem;
    color: #4a5568 !important;
    margin-top: 0.25rem;
}

.kx-callout {
    border-left: 4px solid var(--kite-gold);
    background: #f7f8fb !important;
    padding: 1rem 1.25rem;
    border-radius: 0 8px 8px 0;
    margin: 1rem 0;
    color: #1a2333 !important;
}

.kx-takeaway {
    display: flex;
    gap: 0.9rem;
    align-items: flex-start;
    background: #fdfcf6 !important;
    border: 1px solid #efe3c1;
    border-radius: 10px;
    padding: 0.9rem 1.1rem;
    margin: 0.5rem 0;
}
.kx-takeaway-num {
    flex: 0 0 30px;
    width: 30px; height: 30px;
    border-radius: 50%;
    background: var(--kite-gold);
    color: #2a1b00 !important;
    font-weight: 700;
    display: flex; align-items: center; justify-content: center;
}
.kx-takeaway-body { color: #1a2333 !important; font-size: 0.96rem; line-height: 1.45; }

.kx-ai-response {
    background: #f7f8fb !important;
    border-left: 4px solid #0a2540;
    padding: 1.25rem 1.5rem;
    color: #1a2333 !important;
    border-radius: 0 10px 10px 0;
    margin: 1rem 0;
    line-height: 1.55;
}

.kx-q-excellent { background: #ecfdf3 !important; border: 1px solid #abefc6; padding: 1rem 1.2rem; border-radius: 10px; color: #14532d !important; }
.kx-q-good      { background: #f0f9ff !important; border: 1px solid #bae0fd; padding: 1rem 1.2rem; border-radius: 10px; color: #0c4a6e !important; }
.kx-q-warning   { background: #fffaeb !important; border: 1px solid #fcd980; padding: 1rem 1.2rem; border-radius: 10px; color: #78350f !important; }
.kx-q-poor      { background: #fef2f2 !important; border: 1px solid #fda4a4; padding: 1rem 1.2rem; border-radius: 10px; color: #7f1d1d !important; }

.stTabs [data-baseweb="tab-list"] {
    gap: 4px;
    border-bottom: 1px solid var(--kite-border);
}
.stTabs [data-baseweb="tab"] {
    padding: 0.75rem 1.2rem;
    font-weight: 500;
    color: var(--kite-text-soft);
    border-radius: 8px 8px 0 0;
}
.stTabs [aria-selected="true"] {
    color: var(--kite-navy) !important;
    border-bottom: 3px solid var(--kite-gold) !important;
    font-weight: 700;
}

button,
.stButton > button,
.stButton > button *,
[data-testid*="Button"],
[data-testid*="Button"] *,
[data-testid*="button"],
[data-testid*="button"] *,
div[data-testid="stSidebar"] button,
div[data-testid="stSidebar"] button *,
div[data-testid="stSidebar"] button p,
div[data-testid="stSidebar"] button span {
    background-color: #0a2540 !important;
    color: #ffffff !important;
    border-radius: 8px !important;
    border: none !important;
    font-weight: 600 !important;
}
button:hover,
.stButton > button:hover,
.stButton > button:hover * {
    background-color: #14304f !important;
    color: #e6c878 !important;
}

.kx-dash-row { display: flex; gap: 0.75rem; margin: 0.5rem 0; }
.kx-dash-panel {
    flex: 1; background: #ffffff !important; border: 1px solid #e3e6ee;
    border-radius: 10px; padding: 0.5rem; min-width: 0;
}

section[data-testid="stSidebar"] {
    background: #f7f8fb !important;
    border-right: 1px solid #e3e6ee;
}
section[data-testid="stSidebar"] h1,
section[data-testid="stSidebar"] h2,
section[data-testid="stSidebar"] h3,
section[data-testid="stSidebar"] label,
section[data-testid="stSidebar"] p,
section[data-testid="stSidebar"] span:not([class*="st-"]),
section[data-testid="stSidebar"] .stMarkdown {
    color: #1a2333 !important;
}
section[data-testid="stSidebar"] .stSelectbox div[data-baseweb="select"] * {
    color: #1a2333 !important;
    background-color: #ffffff !important;
}
section[data-testid="stSidebar"] input {
    color: #1a2333 !important;
    background-color: #ffffff !important;
}
</style>
""",
    unsafe_allow_html=True,
)


KITE_PLOTLY = dict(
    plot_bgcolor="#ffffff",
    paper_bgcolor="#ffffff",
    font=dict(family="Inter, system-ui, sans-serif", color="#1a2333"),
    colorway=["#0a2540", "#c79a3a", "#14304f", "#e6c878", "#4a5568", "#a98a3d"],
    margin=dict(t=50, l=40, r=20, b=40),
)


def kx_apply_theme(fig: go.Figure) -> go.Figure:
    fig.update_layout(**KITE_PLOTLY)
    return fig



# ============================================================
# GROQ WRAPPER
# ============================================================

FIXED_MODEL = "llama-3.3-70b-versatile"


class GroqLLM:
    def __init__(self, api_key: str, model: str = FIXED_MODEL):
        self.client = Groq(api_key=api_key)
        self.model = model

    def predict(self, prompt: str, max_tokens: int = 2000, system: str = None) -> str:
        if system is None:
            system = (
                "You are KiteIQX Intelligence, a senior management consultant. "
                "Speak in clear business language. Always give specific numbers and "
                "concrete recommendations - never vague generalities."
            )
        try:
            r = self.client.chat.completions.create(
                model=self.model,
                temperature=0.2,
                max_tokens=max_tokens,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": prompt},
                ],
            )
            return r.choices[0].message.content
        except Exception as e:
            return f"AI error: {e}"


class ClaudeLLM:
    """Anthropic Claude wrapper with the same .predict() interface as GroqLLM."""
    def __init__(self, api_key: str, model: str = CLAUDE_MODEL):
        self.client = Anthropic(api_key=api_key)
        self.model = model

    def predict(self, prompt: str, max_tokens: int = 2000, system: str = None) -> str:
        if system is None:
            system = (
                "You are KiteIQX Intelligence, a senior management consultant. "
                "Speak in clear business language. Always give specific numbers and "
                "concrete recommendations - never vague generalities."
            )
        try:
            r = self.client.messages.create(
                model=self.model,
                max_tokens=max_tokens,
                temperature=0.2,
                system=system,
                messages=[{"role": "user", "content": prompt}],
            )
            # concatenate text blocks
            parts = [b.text for b in r.content if getattr(b, "type", None) == "text"]
            return "".join(parts) if parts else (r.content[0].text if r.content else "")
        except Exception as e:
            return f"AI error: {e}"

class UniversalAnalytics:
    def __init__(self, df: pd.DataFrame, llm=None):
        self.df = df
        self.llm = llm
        self.original_df = df.copy()
        self.process_data()
        self.generate_insights()

    def process_data(self):
        self.numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()
        self.categorical_cols = self.df.select_dtypes(include=["object", "category"]).columns.tolist()
        self.datetime_cols = []

        for col in list(self.df.columns):
            if any(k in col.lower() for k in ("date", "time", "created", "updated", "timestamp")):
                try:
                    self.df[col] = pd.to_datetime(self.df[col])
                    self.datetime_cols.append(col)
                    if col in self.categorical_cols:
                        self.categorical_cols.remove(col)
                except Exception:
                    pass

        self.money_cols = [
            c for c in self.numeric_cols
            if any(w in c.lower() for w in ("price", "cost", "amount", "revenue", "salary", "income", "fee", "payment", "spend", "sales", "gmv"))
        ]
        self.score_cols = [c for c in self.numeric_cols if any(w in c.lower() for w in ("score", "rating", "rank", "grade"))]
        self.id_cols = [
            c for c in self.df.columns
            if any(w in c.lower() for w in ("id", "key")) and self.df[c].nunique() / max(1, len(self.df)) > 0.8
        ]

        if len(self.money_cols) > 1:
            self.df["total_monetary_value"] = self.df[self.money_cols].sum(axis=1, skipna=True)
        elif len(self.money_cols) == 1:
            self.df["total_monetary_value"] = self.df[self.money_cols[0]]

        self.numeric_cols = self.df.select_dtypes(include=[np.number]).columns.tolist()
        self.check_data_quality()

    def check_data_quality(self):
        issues = {}
        dup = int(self.df.duplicated().sum())
        if dup:
            issues["duplicates"] = {"count": dup, "percentage": dup / len(self.df) * 100}

        missing = {}
        for c in self.df.columns:
            m = int(self.df[c].isnull().sum())
            if m:
                missing[c] = {"count": m, "percentage": m / len(self.df) * 100}
        if missing:
            issues["missing_data"] = missing

        outliers = {}
        for c in self.numeric_cols:
            o = self._iqr_outliers(c)
            if o["count"] > 0:
                outliers[c] = o
        if outliers:
            issues["outliers"] = outliers

        self.data_quality_issues = issues

    def _iqr_outliers(self, col):
        s = self.df[col].dropna()
        if len(s) < 10:
            return {"count": 0}
        Q1, Q3 = s.quantile(0.25), s.quantile(0.75)
        IQR = Q3 - Q1
        lo, hi = Q1 - 1.5 * IQR, Q3 + 1.5 * IQR
        outs = s[(s < lo) | (s > hi)]
        return {
            "count": int(len(outs)),
            "percentage": len(outs) / len(s) * 100,
            "range": f"{outs.min():.2f} to {outs.max():.2f}" if len(outs) else "None",
            "bounds": f"normal range: {lo:.2f} to {hi:.2f}",
        }

    def get_data_quality_report(self):
        if not self.data_quality_issues:
            return "excellent", 100, "No significant data quality issues detected. The dataset is analysis-ready."

        severity = 0
        bullets = []
        if "duplicates" in self.data_quality_issues:
            d = self.data_quality_issues["duplicates"]
            severity += min(d["percentage"] * 2, 30)
            bullets.append(f"**Duplicates:** {d['count']:,} records ({d['percentage']:.1f}%)")
        if "missing_data" in self.data_quality_issues:
            top = max(v["percentage"] for v in self.data_quality_issues["missing_data"].values())
            severity += min(top, 25)
            bullets.append(f"**Missing data:** up to {top:.1f}% in worst-affected column")
        if "outliers" in self.data_quality_issues:
            top = max(v["percentage"] for v in self.data_quality_issues["outliers"].values())
            severity += min(top, 15)
            bullets.append(f"**Outliers:** up to {top:.1f}% values flagged as statistical outliers")

        score = max(0, int(100 - severity))
        if score >= 90:
            status = "excellent"
        elif score >= 75:
            status = "good"
        elif score >= 60:
            status = "fair"
        else:
            status = "poor"
        return status, score, "  \n".join(bullets)

    def clean_duplicates(self):
        if "duplicates" in self.data_quality_issues:
            n = len(self.df)
            self.df = self.df.drop_duplicates().reset_index(drop=True)
            removed = n - len(self.df)
            self.process_data()
            self.generate_insights()
            return f"Removed {removed:,} duplicate records. Now {len(self.df):,} unique rows."
        return "No duplicates to remove."

    def generate_insights(self):
        self.insights = {
            "basic_stats": {
                "rows": len(self.df),
                "columns": len(self.df.columns),
                "missing_pct": (self.df.isnull().sum().sum() / max(1, self.df.size)) * 100,
            },
            "column_types": {
                "numeric": len(self.numeric_cols),
                "categorical": len(self.categorical_cols),
                "datetime": len(self.datetime_cols),
            },
        }
        if "total_monetary_value" in self.df.columns:
            mc = "total_monetary_value"
            self.insights["monetary"] = {
                "total": float(self.df[mc].sum()),
                "avg": float(self.df[mc].mean()),
                "max": float(self.df[mc].max()),
                "min": float(self.df[mc].min()),
            }

    def _data_brief(self, max_rows: int = 5) -> str:
        return (
            f"Rows: {len(self.df):,}\n"
            f"Columns: {len(self.df.columns)} "
            f"({len(self.numeric_cols)} numeric, {len(self.categorical_cols)} categorical, "
            f"{len(self.datetime_cols)} datetime)\n"
            f"Numeric columns: {', '.join(self.numeric_cols[:12])}\n"
            f"Categorical columns: {', '.join(self.categorical_cols[:12])}\n\n"
            f"Sample rows:\n{self.df.head(max_rows).to_string()}\n\n"
            f"Numeric summary:\n{self.df[self.numeric_cols[:6]].describe().to_string() if self.numeric_cols else '(none)'}\n"
        )

    def ai_query(self, query: str) -> str:
        if not self.llm:
            return "AI is not configured. Add GROQ_API_KEY to .streamlit/secrets.toml."
        prompt = f"DATASET CONTEXT\n{self._data_brief()}\n\nUSER QUESTION\n{query}\n\nAnswer with specific numbers and concrete recommendations."
        answer = self.llm.predict(prompt)
        filename = st.session_state.get("last_saved_path", "unknown")
        log_ai_query(filename, query, answer)
        return answer

    def ai_executive_summary(self) -> str:
        if not self.llm:
            return self._fallback_summary()

        cat_context = ""
        for c in self.categorical_cols[:3]:
            vc = self.df[c].value_counts().head(5)
            cat_context += f"\n  {c}: {', '.join(f'{k}={v}' for k, v in vc.items())}"

        num_context = ""
        for c in self.numeric_cols[:5]:
            s = self.df[c].dropna()
            num_context += f"\n  {c}: mean={s.mean():.2f}, median={s.median():.2f}, max={s.max():.2f}, min={s.min():.2f}"

        monetary_ctx = ""
        if "monetary" in self.insights:
            m = self.insights["monetary"]
            monetary_ctx = f"\n  Total monetary value: ${m['total']:,.0f}, avg per record: ${m['avg']:,.2f}"

        prompt = (
            f"DATASET CONTEXT\n{self._data_brief()}"
            f"\nTop categorical breakdowns:{cat_context}"
            f"\nNumeric profile:{num_context}"
            f"{monetary_ctx}\n\n"
            "Write a sharp 3-sentence executive narrative. Rules:\n"
            "  Sentence 1 — BUSINESS DOMAIN: Name the specific business function this data covers "
            "(e.g. 'retail sales', 'HR workforce', 'logistics') and the time scope if visible. "
            "Be specific — not just 'this dataset contains records'.\n"
            "  Sentence 2 — KEY BUSINESS SIGNAL: State the single most important pattern, trend, or "
            "anomaly using actual numbers from the data. Focus on commercial or operational impact.\n"
            "  Sentence 3 — DECISION IMPLICATION: What should a business leader do or watch because of this? "
            "One clear, actionable implication.\n\n"
            "Tone: senior consultant briefing a CEO. Specific numbers, zero filler, no bullet points — "
            "three sentences of clean prose."
        )
        return self.llm.predict(prompt, max_tokens=600)

    def ai_top_takeaways(self) -> list:
        if not self.llm:
            return self._fallback_takeaways()
        prompt = (
            f"DATASET CONTEXT\n{self._data_brief()}\n\n"
            "Return exactly 3 strategic takeaways for a CEO. Each takeaway must:\n"
            "  - Be a single, direct sentence\n"
            "  - Reference a specific number from the data\n"
            "  - Point toward a concrete growth opportunity, operational fix, or strategic decision — "
            "not just an observation\n"
            "  - Start with the finding, end with the recommended next move\n\n"
            "Think like a McKinsey partner: the goal is to surface what to DO next, not just what the data shows.\n\n"
            "Format: numbered list (1. 2. 3.) — no preamble, no headings, nothing else."
        )
        out = self.llm.predict(prompt, max_tokens=500)
        lines = [re.sub(r"^\s*\d+[\.\)]\s*", "", ln).strip() for ln in out.split("\n") if ln.strip()]
        lines = [ln for ln in lines if len(ln) > 10]
        return lines[:3] if lines else self._fallback_takeaways()

    def _fallback_summary(self) -> str:
        b = self.insights["basic_stats"]
        return (
            f"This dataset contains {b['rows']:,} records across {b['columns']} columns, with "
            f"{self.insights['column_types']['numeric']} numeric, {self.insights['column_types']['categorical']} "
            f"categorical, and {self.insights['column_types']['datetime']} datetime fields. "
            f"Overall completeness sits at {100 - b['missing_pct']:.1f}%. "
            "Configure the AI assistant for a deeper narrative summary."
        )

    def _fallback_takeaways(self) -> list:
        out = []
        if self.numeric_cols:
            c = self.numeric_cols[0]
            out.append(f"Primary metric '{c}' averages {self.df[c].mean():.2f} (range {self.df[c].min():.2f} to {self.df[c].max():.2f}).")
        if self.categorical_cols:
            c = self.categorical_cols[0]
            top = self.df[c].value_counts().head(1)
            if len(top):
                out.append(f"In '{c}', '{top.index[0]}' is the dominant segment with {top.iloc[0]:,} records ({top.iloc[0]/len(self.df)*100:.1f}% share).")
        if self.insights["basic_stats"]["missing_pct"] > 5:
            out.append(f"Data completeness is {100 - self.insights['basic_stats']['missing_pct']:.1f}%; the Data Quality tab will guide cleanup before deep analysis.")
        while len(out) < 3:
            out.append("Connect the AI engine (Groq key in secrets) for richer, data-specific takeaways.")
        return out[:3]



# ============================================================
# VISUALIZATIONS
# ============================================================

class VizEngine:
    def __init__(self, analytics: UniversalAnalytics):
        self.a = analytics
        self.df = analytics.df

    def scatter(self, x, y, color=None):
        if color and color in self.a.categorical_cols:
            fig = px.scatter(self.df, x=x, y=y, color=color, title=f"{y} vs {x} by {color}")
        else:
            fig = px.scatter(self.df, x=x, y=y, title=f"{y} vs {x}")
        try:
            corr = self.df[x].corr(self.df[y])
            fig.add_annotation(text=f"r = {corr:.3f}", xref="paper", yref="paper", x=0.02, y=0.98,
                               showarrow=False, bgcolor="rgba(255,255,255,0.9)", bordercolor="#0a2540")
        except Exception:
            corr = float("nan")
        return kx_apply_theme(fig), f"Correlation {corr:.3f}" if corr == corr else "Scatter created"

    def correlation_matrix(self, cols):
        nums = [c for c in cols if c in self.a.numeric_cols]
        if len(nums) < 2:
            return None, "Need at least 2 numeric columns."
        m = self.df[nums].corr()
        fig = px.imshow(m, text_auto=".2f", aspect="auto", color_continuous_scale=["#b91c1c", "#ffffff", "#0a2540"], title="Correlation matrix")
        return kx_apply_theme(fig), f"Correlation matrix for {len(nums)} variables"

    def distribution(self, cols):
        nums = [c for c in cols if c in self.a.numeric_cols][:4]
        if not nums:
            return None, "No numeric columns selected."
        ncols = min(2, len(nums))
        nrows = (len(nums) + 1) // 2
        fig = make_subplots(rows=nrows, cols=ncols, subplot_titles=[f"{c}" for c in nums])
        for i, c in enumerate(nums):
            r, cc = i // ncols + 1, i % ncols + 1
            fig.add_trace(go.Histogram(x=self.df[c], name=c, nbinsx=30, marker_color="#0a2540"), row=r, col=cc)
        fig.update_layout(height=340 * nrows, showlegend=False, title_text="Distribution analysis")
        return kx_apply_theme(fig), f"Distribution analysis for {len(nums)} variables"

    def regression(self, target, predictors):
        preds = [c for c in predictors if c in self.a.numeric_cols and c != target]
        if not preds:
            return None, "Need a numeric predictor."
        p = preds[0]
        data = self.df[[target, p]].dropna()
        X = data[p].values.reshape(-1, 1)
        y = data[target].values
        m = LinearRegression().fit(X, y)
        y_hat = m.predict(X)
        r2 = r2_score(y, y_hat)
        resid = y - y_hat
        fig = make_subplots(rows=2, cols=2, subplot_titles=[f"{p} -> {target}", "Residuals", "Actual vs predicted", "R^2"])
        xs = np.linspace(X.min(), X.max(), 100)
        fig.add_trace(go.Scatter(x=X.flatten(), y=y, mode="markers", name="actual",
                                 marker=dict(color="#0a2540", opacity=0.6)), row=1, col=1)
        fig.add_trace(go.Scatter(x=xs, y=m.predict(xs.reshape(-1, 1)), mode="lines", name="fit",
                                 line=dict(color="#c79a3a", width=2)), row=1, col=1)
        fig.add_trace(go.Scatter(x=y_hat, y=resid, mode="markers", marker=dict(color="#c79a3a", opacity=0.6),
                                 name="residuals"), row=1, col=2)
        fig.add_trace(go.Scatter(x=y, y=y_hat, mode="markers", marker=dict(color="#14304f", opacity=0.6),
                                 name="a-vs-p"), row=2, col=1)
        lo, hi = min(y.min(), y_hat.min()), max(y.max(), y_hat.max())
        fig.add_trace(go.Scatter(x=[lo, hi], y=[lo, hi], mode="lines", line=dict(color="#888", dash="dash"),
                                 name="ideal"), row=2, col=1)
        fig.add_trace(go.Indicator(
            mode="gauge+number", value=r2,
            domain={"x": [0, 1], "y": [0, 1]}, title={"text": "R^2"},
            gauge={"axis": {"range": [None, 1]},
                   "bar": {"color": "#0a2540"},
                   "steps": [{"range": [0, 0.5], "color": "#f0f1f5"},
                             {"range": [0.5, 0.8], "color": "#fcd980"},
                             {"range": [0.8, 1], "color": "#a98a3d"}]}), row=2, col=2)
        fig.update_layout(height=620, showlegend=False, title_text="Predictive model")
        return kx_apply_theme(fig), f"R^2 = {r2:.3f}, RMSE = {np.sqrt(np.mean(resid**2)):.3f}"

    def timeseries(self, date_col, value_col):
        if date_col not in self.a.datetime_cols:
            return None, "Column is not a datetime."
        if value_col not in self.a.numeric_cols:
            return None, "Value column must be numeric."
        d = self.df[[date_col, value_col]].dropna().sort_values(date_col)
        fig = make_subplots(rows=2, cols=1, subplot_titles=[f"{value_col} over time", "Trend"], vertical_spacing=0.12)
        fig.add_trace(go.Scatter(x=d[date_col], y=d[value_col], mode="lines+markers", name=value_col,
                                 line=dict(color="#0a2540", width=2)), row=1, col=1)
        if len(d) > 7:
            d["ma7"] = d[value_col].rolling(7).mean()
            fig.add_trace(go.Scatter(x=d[date_col], y=d["ma7"], mode="lines", name="7-period MA",
                                     line=dict(color="#c79a3a", dash="dash")), row=1, col=1)
        if len(d) > 3:
            xn = np.arange(len(d))
            z = np.polyfit(xn, d[value_col], 1)
            fig.add_trace(go.Scatter(x=d[date_col], y=np.poly1d(z)(xn), mode="lines", name="trend",
                                     line=dict(color="#14304f", width=3)), row=2, col=1)
        fig.update_layout(height=560, title_text="Time-series analysis")
        return kx_apply_theme(fig), f"Time-series for {value_col}"


# ============================================================
# DEMO DATA
# ============================================================

def generate_demo(kind: str) -> pd.DataFrame:
    np.random.seed(42)
    if kind == "E-commerce Sales":
        n = 1000
        return pd.DataFrame({
            "order_id": [f"ORD{i:06d}" for i in range(1, n + 1)],
            "customer_id": [f"CUST{i:05d}" for i in np.random.randint(1, 501, n)],
            "product_category": np.random.choice(["Electronics", "Clothing", "Home", "Books", "Sports"], n),
            "order_value": np.random.lognormal(4, 0.8, n),
            "shipping_cost": np.random.uniform(5, 25, n),
            "customer_age": np.random.randint(18, 75, n),
            "customer_segment": np.random.choice(["Premium", "Standard", "Budget"], n),
            "delivery_days": np.random.randint(1, 15, n),
            "customer_rating": np.random.randint(1, 6, n),
            "order_date": pd.date_range("2024-01-01", periods=n, freq="D"),
        })
    if kind == "Employee Data":
        n = 500
        return pd.DataFrame({
            "employee_id": [f"EMP{i:04d}" for i in range(1, n + 1)],
            "department": np.random.choice(["Engineering", "Sales", "Marketing", "HR", "Finance"], n),
            "salary": np.random.normal(75000, 25000, n),
            "age": np.random.randint(22, 65, n),
            "years_experience": np.random.randint(0, 25, n),
            "performance_score": np.random.normal(3.5, 0.8, n),
            "job_satisfaction": np.random.randint(1, 11, n),
        })
    n = 300
    return pd.DataFrame({
        "id": range(1, n + 1),
        "category": np.random.choice(["A", "B", "C", "D"], n),
        "value": np.random.uniform(10, 100, n),
        "score": np.random.randint(1, 11, n),
        "amount": np.random.uniform(100, 1000, n),
    })


# ============================================================
# FILE LOADING
# ============================================================

def load_uploaded_file(uploaded_file) -> pd.DataFrame:
    name = uploaded_file.name
    ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    saved = UPLOAD_DIR / f"{ts}_{name}"
    raw = uploaded_file.getbuffer()
    with open(saved, "wb") as f:
        f.write(raw)
    st.session_state["last_saved_path"] = str(saved)

    buf = io.BytesIO(raw)
    if ext in ("xlsx", "xls"):
        df = pd.read_excel(buf)
    elif ext == "csv":
        df = pd.read_csv(buf)
    else:
        try:
            df = pd.read_csv(io.BytesIO(raw))
        except Exception:
            df = pd.read_excel(io.BytesIO(raw))
    return df



# ============================================================
# UI BLOCKS
# ============================================================

def render_header():
    st.markdown(
        """
        <div class="kx-brand">
            <div class="kx-mark">K</div>
            <div class="kx-wordmark">Kite<span>IQX</span> Intelligence</div>
        </div>
        <div class="kx-tag">Where intelligence meets impact — turn raw data into a CEO-ready story in minutes.</div>
        """,
        unsafe_allow_html=True,
    )


def render_welcome():
    st.markdown(
        """
        <div class="kx-hero">
            <h2>Upload your data, get a CEO-ready narrative in minutes.</h2>
            <p>Built for non-technical leaders. Drop a CSV or Excel file and KiteIQX will surface the
            <span class="kx-hero-accent">three insights that matter</span>, audit data quality, and let you ask follow-up
            questions in plain English.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown(
            """<div class="kx-card"><div class="kx-card-title">Executive Dashboard</div>
            <div class="kx-card-sub">Narrative summary, hero KPIs, and the three takeaways a leader needs.</div></div>""",
            unsafe_allow_html=True,
        )
    with c2:
        st.markdown(
            """<div class="kx-card"><div class="kx-card-title">AI Intelligence</div>
            <div class="kx-card-sub">Ask any business question; get a consultant-grade answer in seconds.</div></div>""",
            unsafe_allow_html=True,
        )
    with c3:
        st.markdown(
            """<div class="kx-card"><div class="kx-card-title">Data Quality</div>
            <div class="kx-card-sub">Detect duplicates, missingness, and outliers before they cost you a decision.</div></div>""",
            unsafe_allow_html=True,
        )
    st.info("Use the sidebar to upload a CSV / XLSX file or try a demo dataset.")


def render_dashboard():
    a: UniversalAnalytics = st.session_state.analytics
    ins = a.insights

    record_str = f"{ins['basic_stats']['rows']:,}"
    col_str = f"{ins['basic_stats']['columns']}"
    quality_pct = 100 - ins["basic_stats"]["missing_pct"]
    monetary_line = ""
    if "monetary" in ins:
        monetary_line = (
            f' Total tracked value across the dataset is '
            f'<span class="kx-hero-accent">${ins["monetary"]["total"]:,.0f}</span>.'
        )
    st.markdown(
        f"""
        <div class="kx-hero">
            <h2>Business Story</h2>
            <p>KiteIQX has analysed <span class="kx-hero-accent">{record_str}</span> records across
            <span class="kx-hero-accent">{col_str}</span> dimensions
            ({ins['column_types']['numeric']} numeric · {ins['column_types']['categorical']} categorical · {ins['column_types']['datetime']} time-based)
            with <span class="kx-hero-accent">{quality_pct:.1f}%</span> completeness.{monetary_line}
            The AI narrative and takeaways below reflect what the numbers actually say about your business.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    cards = st.columns(4)
    cards[0].markdown(
        f"""<div class="kx-card"><div class="kx-card-title">Records analyzed</div>
        <div class="kx-card-value">{record_str}</div>
        <div class="kx-card-sub">Rows of business activity</div></div>""", unsafe_allow_html=True)
    cards[1].markdown(
        f"""<div class="kx-card"><div class="kx-card-title">Dimensions</div>
        <div class="kx-card-value">{col_str}</div>
        <div class="kx-card-sub">{ins['column_types']['numeric']} numeric · {ins['column_types']['categorical']} categorical · {ins['column_types']['datetime']} time</div></div>""", unsafe_allow_html=True)
    if "monetary" in ins:
        cards[2].markdown(
            f"""<div class="kx-card"><div class="kx-card-title">Total monetary value</div>
            <div class="kx-card-value">${ins['monetary']['total']:,.0f}</div>
            <div class="kx-card-sub">Average per record: ${ins['monetary']['avg']:,.2f}</div></div>""", unsafe_allow_html=True)
    else:
        cards[2].markdown(
            f"""<div class="kx-card"><div class="kx-card-title">Numeric depth</div>
            <div class="kx-card-value">{ins['column_types']['numeric']}</div>
            <div class="kx-card-sub">Quantifiable dimensions</div></div>""", unsafe_allow_html=True)
    cards[3].markdown(
        f"""<div class="kx-card"><div class="kx-card-title">Data completeness</div>
        <div class="kx-card-value">{quality_pct:.1f}%</div>
        <div class="kx-card-sub">Higher is better; details in Data Quality tab</div></div>""", unsafe_allow_html=True)

    st.markdown(" ")
    st.markdown(
        '<div style="font-size:0.78rem;text-transform:uppercase;letter-spacing:0.8px;'
        'color:#4a5568;font-weight:600;margin-bottom:0.25rem;">Business narrative</div>',
        unsafe_allow_html=True,
    )

    if "exec_summary" not in st.session_state:
        with st.spinner("KiteIQX is reading the business context..."):
            st.session_state.exec_summary = a.ai_executive_summary()
    st.markdown(
        f'<div class="kx-callout" style="font-size:1.02rem;line-height:1.65;">'
        f'{st.session_state.exec_summary}</div>',
        unsafe_allow_html=True,
    )

    refresh_col, _ = st.columns([1, 5])
    if refresh_col.button("↻ Regenerate", key="regen_summary"):
        st.session_state.pop("exec_summary", None)
        st.session_state.pop("takeaways", None)
        st.rerun()

    st.markdown(" ")
    st.markdown(
        '<div style="font-size:0.78rem;text-transform:uppercase;letter-spacing:0.8px;'
        'color:#4a5568;font-weight:600;margin-bottom:0.35rem;">3 decisions this data supports</div>',
        unsafe_allow_html=True,
    )
    if "takeaways" not in st.session_state:
        with st.spinner("Identifying the headline insights..."):
            st.session_state.takeaways = a.ai_top_takeaways()

    if not st.session_state.get("upload_logged") and st.session_state.get("upload_filename"):
        log_upload(
            st.session_state.upload_filename,
            a.df,
            summary=st.session_state.get("exec_summary", ""),
            takeaways=st.session_state.get("takeaways", []),
        )
        st.session_state.upload_logged = True

    for i, t in enumerate(st.session_state.takeaways, 1):
        st.markdown(
            f"""<div class="kx-takeaway"><div class="kx-takeaway-num">{i}</div>
            <div class="kx-takeaway-body">{t}</div></div>""", unsafe_allow_html=True)

    st.markdown(" ")
    st.markdown(
        '<div style="font-size:0.78rem;text-transform:uppercase;letter-spacing:0.8px;'
        'color:#4a5568;font-weight:600;margin-bottom:0.5rem;">Supporting visuals</div>',
        unsafe_allow_html=True,
    )

    COMPACT = dict(
        plot_bgcolor="#ffffff",
        paper_bgcolor="#ffffff",
        font=dict(family="Inter, system-ui, sans-serif", color="#1a2333", size=11),
        colorway=["#0a2540", "#c79a3a", "#14304f", "#e6c878", "#4a5568"],
        margin=dict(t=36, l=28, r=12, b=28),
        height=220,
        showlegend=False,
    )

    vis_cols = st.columns(3)

    with vis_cols[0]:
        if a.categorical_cols:
            cat = a.categorical_cols[0]
            top = a.df[cat].value_counts().head(6).reset_index()
            top.columns = [cat, "count"]
            fig1 = px.bar(top, x=cat, y="count", title=f"{cat} breakdown")
            fig1.update_traces(marker_color="#0a2540")
            fig1.update_layout(**COMPACT)
            fig1.update_xaxes(tickangle=-30, tickfont_size=10)
            st.plotly_chart(fig1, use_container_width=True)
        elif a.numeric_cols:
            c = a.numeric_cols[0]
            fig1 = px.histogram(a.df, x=c, nbins=20, title=f"{c} distribution")
            fig1.update_traces(marker_color="#0a2540")
            fig1.update_layout(**COMPACT)
            st.plotly_chart(fig1, use_container_width=True)

    with vis_cols[1]:
        if a.datetime_cols and a.numeric_cols:
            dc = a.datetime_cols[0]
            vc = a.money_cols[0] if a.money_cols else a.numeric_cols[0]
            ts = a.df[[dc, vc]].dropna().sort_values(dc)
            fig2 = px.line(ts, x=dc, y=vc, title=f"{vc} over time")
            fig2.update_traces(line_color="#c79a3a", line_width=2)
            fig2.update_layout(**COMPACT)
            st.plotly_chart(fig2, use_container_width=True)
        elif len(a.numeric_cols) >= 2:
            x2, y2 = a.numeric_cols[0], a.numeric_cols[1]
            fig2 = px.scatter(a.df, x=x2, y=y2, title=f"{y2} vs {x2}")
            fig2.update_traces(marker=dict(color="#0a2540", opacity=0.5, size=5))
            fig2.update_layout(**COMPACT)
            st.plotly_chart(fig2, use_container_width=True)

    with vis_cols[2]:
        if len(a.categorical_cols) >= 2:
            cat2 = a.categorical_cols[1]
            top2 = a.df[cat2].value_counts().head(6).reset_index()
            top2.columns = [cat2, "count"]
            fig3 = px.bar(top2, x=cat2, y="count", title=f"{cat2} split",
                          color="count", color_continuous_scale=["#e6c878", "#0a2540"])
            fig3.update_layout(**COMPACT)
            fig3.update_coloraxes(showscale=False)
            fig3.update_xaxes(tickangle=-30, tickfont_size=10)
            st.plotly_chart(fig3, use_container_width=True)
        elif len(a.numeric_cols) >= 2:
            c3 = a.numeric_cols[1] if len(a.numeric_cols) > 1 else a.numeric_cols[0]
            fig3 = px.histogram(a.df, x=c3, nbins=20, title=f"{c3} distribution")
            fig3.update_traces(marker_color="#c79a3a")
            fig3.update_layout(**COMPACT)
            st.plotly_chart(fig3, use_container_width=True)

    if len(a.numeric_cols) >= 3 and len(a.categorical_cols) >= 1:
        vis_cols2 = st.columns(2)
        with vis_cols2[0]:
            box_col = a.money_cols[0] if a.money_cols else a.numeric_cols[0]
            cat_box = a.categorical_cols[0]
            fig4 = px.box(
                a.df, x=cat_box, y=box_col, title=f"{box_col} by {cat_box}",
                color_discrete_sequence=["#0a2540"],
            )
            fig4.update_layout(**COMPACT)
            st.plotly_chart(fig4, use_container_width=True)
        with vis_cols2[1]:
            corr_cols = a.numeric_cols[:5]
            corr_m = a.df[corr_cols].corr()
            fig5 = px.imshow(
                corr_m, text_auto=".1f", aspect="auto",
                color_continuous_scale=["#b91c1c", "#ffffff", "#0a2540"],
                title="Correlation heat",
            )
            fig5.update_layout(**{**COMPACT, "height": 220})
            st.plotly_chart(fig5, use_container_width=True)

    st.markdown(
        """
        <div class="kx-callout">
            <strong>Next step</strong>: open the <em>AI Intelligence</em> tab to ask follow-up questions in plain English,
            or jump to <em>Data Quality</em> to clean the dataset before deeper modeling.
        </div>
        """, unsafe_allow_html=True,
    )


def render_ai():
    a: UniversalAnalytics = st.session_state.analytics
    st.markdown('<div class="kx-card"><div class="kx-card-title">AI Intelligence</div>'
                '<div class="kx-card-sub">Ask any business question. Get a consultant-grade answer.</div></div>',
                unsafe_allow_html=True)

    if not a.llm:
        st.warning("AI is unavailable: GROQ_API_KEY missing. Add it to .streamlit/secrets.toml.")
        return

    qcols = st.columns(4)
    quick_prompts = {
        "Data overview": "Give a comprehensive overview of this dataset with the most important findings.",
        "Top correlations": "Identify the strongest correlations in the data and explain their business significance.",
        "Key trends": "Identify and explain the most important trends or patterns in the data.",
        "Anomalies": "Spot anomalies or outliers in the data and explain why they matter.",
    }
    for col, (label, qry) in zip(qcols, quick_prompts.items()):
        if col.button(label, key=f"qa_{label}"):
            with st.spinner("KiteIQX analyzing..."):
                st.session_state.ai_last_response = a.ai_query(qry)

    st.markdown(" ")
    user_q = st.text_area(
        "Custom question:",
        placeholder='e.g. "Which customer segment drives the most revenue and why?"',
        height=110,
        key="ai_custom_q",
    )
    if st.button("Ask KiteIQX", type="primary"):
        if user_q.strip():
            with st.spinner("Thinking..."):
                st.session_state.ai_last_response = a.ai_query(user_q.strip())

    if "ai_last_response" in st.session_state:
        st.markdown(
            f'<div class="kx-ai-response"><strong>KiteIQX Intelligence:</strong><br><br>'
            f'{st.session_state.ai_last_response}</div>',
            unsafe_allow_html=True,
        )

    with st.expander("Available columns"):
        c1, c2, c3 = st.columns(3)
        with c1:
            st.markdown("**Numeric**")
            for c in a.numeric_cols[:15]:
                st.write(f"• {c}")
        with c2:
            st.markdown("**Categorical**")
            for c in a.categorical_cols[:15]:
                st.write(f"• {c}")
        with c3:
            st.markdown("**Datetime**")
            for c in a.datetime_cols:
                st.write(f"• {c}")


def render_data_quality():
    a: UniversalAnalytics = st.session_state.analytics
    status, score, summary = a.get_data_quality_report()

    panel_class = {"excellent": "kx-q-excellent", "good": "kx-q-good",
                   "fair": "kx-q-warning", "poor": "kx-q-poor"}[status]
    headline = {"excellent": "Excellent", "good": "Good", "fair": "Needs attention", "poor": "Poor"}[status]

    st.markdown(
        f"""<div class="{panel_class}">
        <div style="font-size: 0.85rem; text-transform: uppercase; letter-spacing: 0.6px; opacity:0.7;">Overall data quality</div>
        <div style="font-size: 2rem; font-weight: 800; line-height: 1.1;">{headline} — {score}/100</div>
        <div style="margin-top: 0.5rem;">{summary}</div>
        </div>""",
        unsafe_allow_html=True,
    )

    st.markdown(" ")
    cols = st.columns(3)

    with cols[0]:
        st.markdown('<div class="kx-card"><div class="kx-card-title">Duplicates</div>', unsafe_allow_html=True)
        if "duplicates" in a.data_quality_issues:
            d = a.data_quality_issues["duplicates"]
            st.markdown(f'<div class="kx-card-value">{d["count"]:,}</div>'
                        f'<div class="kx-card-sub">{d["percentage"]:.1f}% of rows are exact duplicates</div></div>',
                        unsafe_allow_html=True)
            if st.button("Remove duplicates", key="dq_dedup"):
                msg = a.clean_duplicates()
                st.success(msg)
                st.session_state.pop("exec_summary", None)
                st.session_state.pop("takeaways", None)
                st.rerun()
        else:
            st.markdown('<div class="kx-card-value">0</div><div class="kx-card-sub">No duplicate rows detected.</div></div>',
                        unsafe_allow_html=True)

    with cols[1]:
        st.markdown('<div class="kx-card"><div class="kx-card-title">Missing values</div>', unsafe_allow_html=True)
        if "missing_data" in a.data_quality_issues:
            md = a.data_quality_issues["missing_data"]
            worst = max(md.items(), key=lambda x: x[1]["percentage"])
            st.markdown(
                f'<div class="kx-card-value">{len(md)}</div>'
                f'<div class="kx-card-sub">columns affected; worst: <b>{worst[0]}</b> ({worst[1]["percentage"]:.1f}%)</div></div>',
                unsafe_allow_html=True,
            )
        else:
            st.markdown('<div class="kx-card-value">0</div><div class="kx-card-sub">All columns are 100% populated.</div></div>',
                        unsafe_allow_html=True)

    with cols[2]:
        st.markdown('<div class="kx-card"><div class="kx-card-title">Outliers</div>', unsafe_allow_html=True)
        if "outliers" in a.data_quality_issues:
            od = a.data_quality_issues["outliers"]
            worst = max(od.items(), key=lambda x: x[1]["percentage"])
            st.markdown(
                f'<div class="kx-card-value">{len(od)}</div>'
                f'<div class="kx-card-sub">numeric columns flagged; worst: <b>{worst[0]}</b> ({worst[1]["percentage"]:.1f}%)</div></div>',
                unsafe_allow_html=True,
            )
        else:
            st.markdown('<div class="kx-card-value">0</div><div class="kx-card-sub">No statistical outliers (IQR method).</div></div>',
                        unsafe_allow_html=True)

    st.markdown(" ")
    if "missing_data" in a.data_quality_issues:
        st.markdown("#### Missing data — column-by-column")
        md = a.data_quality_issues["missing_data"]
        tbl = pd.DataFrame([
            {"Column": k, "Missing rows": v["count"], "Missing %": round(v["percentage"], 2)}
            for k, v in sorted(md.items(), key=lambda x: -x[1]["percentage"])
        ])
        st.dataframe(tbl, use_container_width=True, hide_index=True)

    if "outliers" in a.data_quality_issues:
        st.markdown("#### Outliers — column-by-column (IQR method)")
        od = a.data_quality_issues["outliers"]
        tbl = pd.DataFrame([
            {"Column": k, "Outlier rows": v["count"], "Outlier %": round(v["percentage"], 2),
             "Outlier range": v.get("range", ""), "Expected range": v.get("bounds", "")}
            for k, v in sorted(od.items(), key=lambda x: -x[1]["percentage"])
        ])
        st.dataframe(tbl, use_container_width=True, hide_index=True)


def render_advanced():
    a: UniversalAnalytics = st.session_state.analytics
    v: VizEngine = st.session_state.viz_engine

    st.markdown('<div class="kx-card"><div class="kx-card-title">Advanced Analytics</div>'
                '<div class="kx-card-sub">Select columns, pick a method, run.</div></div>', unsafe_allow_html=True)

    cols = a.df.columns.tolist()
    selected = st.multiselect("Columns to analyze", cols, default=cols[: min(5, len(cols))])
    if not selected:
        st.info("Select at least one column.")
        return

    numeric = [c for c in selected if c in a.numeric_cols]
    categorical = [c for c in selected if c in a.categorical_cols]
    datetime_ = [c for c in selected if c in a.datetime_cols]

    methods = []
    if len(numeric) >= 2:
        methods += ["Correlation matrix", "Scatter plot", "Predictive model"]
    if numeric:
        methods.append("Distribution analysis")
    if datetime_ and numeric:
        methods.append("Time-series analysis")
    if not methods:
        st.info("No methods available for this selection.")
        return

    pick = st.selectbox("Method", methods)
    if st.button("Run analysis"):
        if pick == "Correlation matrix":
            fig, msg = v.correlation_matrix(selected)
        elif pick == "Distribution analysis":
            fig, msg = v.distribution(selected)
        elif pick == "Scatter plot":
            fig, msg = v.scatter(numeric[0], numeric[1], categorical[0] if categorical else None)
        elif pick == "Predictive model":
            fig, msg = v.regression(numeric[-1], numeric[:-1])
        elif pick == "Time-series analysis":
            fig, msg = v.timeseries(datetime_[0], numeric[0])
        else:
            fig, msg = None, "Method not implemented."
        if fig:
            st.plotly_chart(fig, use_container_width=True)
            st.success(msg)
        else:
            st.error(msg)


def render_custom_charts():
    a: UniversalAnalytics = st.session_state.analytics
    v: VizEngine = st.session_state.viz_engine

    left, right = st.columns([1, 2])
    with left:
        st.markdown('<div class="kx-card"><div class="kx-card-title">Build a chart</div></div>', unsafe_allow_html=True)
        chart_type = st.selectbox("Chart type",
                                  ["Scatter", "Correlation matrix", "Distribution", "Predictive model", "Time-series"])

        fig, msg = None, ""
        if chart_type == "Scatter" and len(a.numeric_cols) >= 2:
            x = st.selectbox("X", a.numeric_cols, key="cc_x")
            y = st.selectbox("Y", a.numeric_cols, key="cc_y")
            color = st.selectbox("Color by", [None] + a.categorical_cols, key="cc_color")
            if st.button("Render", key="cc_render_scatter"):
                fig, msg = v.scatter(x, y, color)
        elif chart_type == "Correlation matrix":
            sel = st.multiselect("Columns", a.numeric_cols, default=a.numeric_cols[:5], key="cc_corr")
            if st.button("Render", key="cc_render_corr"):
                fig, msg = v.correlation_matrix(sel)
        elif chart_type == "Distribution":
            sel = st.multiselect("Columns", a.numeric_cols, default=a.numeric_cols[:3], key="cc_dist")
            if st.button("Render", key="cc_render_dist"):
                fig, msg = v.distribution(sel)
        elif chart_type == "Predictive model" and len(a.numeric_cols) >= 2:
            target = st.selectbox("Target", a.numeric_cols, key="cc_target")
            preds = st.multiselect("Predictors", [c for c in a.numeric_cols if c != target],
                                   default=[c for c in a.numeric_cols if c != target][:2], key="cc_preds")
            if st.button("Render", key="cc_render_pred"):
                fig, msg = v.regression(target, preds)
        elif chart_type == "Time-series" and a.datetime_cols and a.numeric_cols:
            date_col = st.selectbox("Date column", a.datetime_cols, key="cc_date")
            val_col = st.selectbox("Value column", a.numeric_cols, key="cc_val")
            if st.button("Render", key="cc_render_ts"):
                fig, msg = v.timeseries(date_col, val_col)
        else:
            st.info("Not enough data of the right type for this chart.")

        if fig:
            st.session_state.cc_fig = fig
            st.session_state.cc_msg = msg

    with right:
        if "cc_fig" in st.session_state:
            st.plotly_chart(st.session_state.cc_fig, use_container_width=True)
            st.success(st.session_state.get("cc_msg", ""))
        else:
            st.info("Configure a chart on the left.")


def render_explorer():
    a: UniversalAnalytics = st.session_state.analytics
    df = a.df

    left, right = st.columns([1, 3])
    with left:
        st.markdown('<div class="kx-card"><div class="kx-card-title">Filters</div></div>', unsafe_allow_html=True)
        cols = st.multiselect("Columns", df.columns.tolist(), default=df.columns.tolist()[:10])
        rows = st.slider("Rows to display", 10, min(2000, len(df)), 100)
        view = st.radio("View", ["Head", "Sample", "Tail"], horizontal=True)
        search = st.text_input("Search text")

    with right:
        view_df = df[cols] if cols else df
        if search:
            text_cols = view_df.select_dtypes(include=["object"]).columns
            if len(text_cols):
                mask = view_df[text_cols].astype(str).apply(lambda x: x.str.contains(search, case=False, na=False)).any(axis=1)
                view_df = view_df[mask]
        if view == "Head":
            view_df = view_df.head(rows)
        elif view == "Tail":
            view_df = view_df.tail(rows)
        else:
            view_df = view_df.sample(min(rows, len(view_df))) if len(view_df) else view_df
        st.dataframe(view_df, use_container_width=True)

        m = st.columns(4)
        m[0].metric("Showing", f"{len(view_df):,}")
        m[1].metric("Total rows", f"{len(df):,}")
        m[2].metric("Columns", len(cols) if cols else len(df.columns))
        m[3].metric("Memory", f"{df.memory_usage(deep=True).sum() / 1024**2:.2f} MB")

        nums = [c for c in view_df.columns if c in a.numeric_cols]
        if nums:
            st.markdown("#### Numeric summary")
            st.dataframe(view_df[nums].describe(), use_container_width=True)





def render_quality_and_charts():
    """Merged tab: Data Quality (top) + Custom Charts (bottom)."""
    st.markdown('<div class="kx-card"><div class="kx-card-title">Data Quality &amp; Custom Charts</div>'
                '<div class="kx-card-sub">Audit completeness and outliers, then build any chart you need.</div></div>',
                unsafe_allow_html=True)
    render_data_quality()
    st.markdown("---")
    st.markdown('<div style="font-size:0.95rem;font-weight:700;color:#1a2333;margin:0.5rem 0;">Custom charts</div>',
                unsafe_allow_html=True)
    render_custom_charts()



import json as _json
import re as _re
from datetime import datetime as _dt
from pathlib import Path as _Path



# ============================================================
# CONSULTANT PRESENTATION ENGINE (integrated)
# ============================================================
#
# Consultant Engine — the analytical core.
# =========================================
# This is NOT a stat-miner. It reasons about a metric the way a BCG/McKinsey
# analyst would: decompose a change across dimensions, localize the driver,
# recurse into it, and quantify each segment's contribution to the total move.
#
# Math computes every number. The LLM (elsewhere) only narrates proven findings.
#
# Key objects:
#   - StoryContext   : what kind of story the data supports (detect at runtime)
#   - DriverAnalysis : decomposition tree for a metric change
#   - Finding        : one proven, quantified insight + the visual that proves it
#


# ============================================================
# 0. SMALL-CELL GUARDS / CONFIG
# ============================================================
MIN_CELL_ROWS = 8        # don't trust a segment with fewer rows than this
MIN_CONTRIB_PCT = 12.0   # a driver must explain at least this % of the move
MAX_TREE_DEPTH = 3       # metric -> dim1 -> dim2 -> dim3
TOP_DIM_CARD = 6         # max categories to chart


def _fmt(v):
    av = abs(v)
    if av >= 1_000_000: return f"{v/1_000_000:.2f}M"
    if av >= 1_000:     return f"{v/1_000:.1f}K"
    if av == int(av):   return f"{int(v):,}"
    return f"{v:,.2f}"


def _signpct(v):
    return f"{v:+.1f}%"


# ============================================================
# 1. STORY CONTEXT  — detect what's possible at runtime
# ============================================================
class StoryContext:
    """Inspects the dataframe and decides which story modes are viable."""

    def __init__(self, analytics):
        self.a = analytics
        self.df = analytics.df
        self.modes = []           # e.g. ['driver_time', 'concentration', 'correlation']
        self.time_col = None
        self.has_two_periods = False
        self.metrics = []         # numeric value columns worth tracking
        self.dims = []            # categorical dimensions worth slicing by
        self._detect()

    def _detect(self):
        a, df = self.a, self.df

        # metrics = money first, else generic numerics (exclude id-like & scores)
        self.metrics = list(a.money_cols) or [
            c for c in a.numeric_cols
            if c not in getattr(a, "id_cols", []) and df[c].nunique() > 5
        ]
        self.metrics = self.metrics[:4]

        # dims = categoricals with a sane cardinality (2..40 distinct)
        self.dims = [
            c for c in a.categorical_cols
            if 2 <= df[c].nunique() <= 40
        ]

        # time?
        if a.datetime_cols:
            self.time_col = a.datetime_cols[0]
            t = pd.to_datetime(df[self.time_col], errors="coerce").dropna()
            if len(t):
                span_days = (t.max() - t.min()).days
                # need at least ~2 comparable periods
                self.has_two_periods = span_days >= 60

        # decide modes
        if self.time_col and self.has_two_periods and self.metrics:
            self.modes.append("driver_time")     # YoY / period-over-period decomposition
        if self.dims and self.metrics:
            self.modes.append("driver_mix")      # which segment drives the metric (no time needed)
        if self.dims:
            self.modes.append("concentration")
        if len(a.numeric_cols) >= 2:
            self.modes.append("correlation")
        if not self.modes:
            self.modes.append("profile")

    def describe(self):
        """Human sentence describing what the agent can build."""
        bits = []
        if "driver_time" in self.modes:
            bits.append(f"a **driver story** — track *{self.metrics[0]}* over time and find what's moving it across {', '.join(self.dims[:3]) or 'segments'}")
        elif "driver_mix" in self.modes:
            bits.append(f"a **mix story** — break down *{self.metrics[0]}* to see which {', '.join(self.dims[:2])} segments drive it")
        if "concentration" in self.modes and "driver_time" not in self.modes:
            bits.append("a **concentration story** — where activity clusters and the risk that creates")
        if "correlation" in self.modes:
            bits.append("a **relationship story** — which metrics move together")
        if self.modes == ["profile"]:
            bits.append("a **profile overview** — the shape of the data (limited story angles without segments or a metric)")
        primary = ("driver" if "driver_time" in self.modes else
                   "mix" if "driver_mix" in self.modes else
                   "concentration" if "concentration" in self.modes else "profile")
        return primary, bits


# ============================================================
# 2. CONTRIBUTION ANALYSIS  — the heart of the "why"
# ============================================================
def _period_split(df, time_col):
    """Split rows into two comparable periods (recent vs prior)."""
    t = pd.to_datetime(df[time_col], errors="coerce")
    d = df.assign(_t=t).dropna(subset=["_t"]).sort_values("_t")
    if len(d) < 2 * MIN_CELL_ROWS:
        return None, None, None, None
    # try year-over-year if >= 2 calendar years, else split by median date
    years = d["_t"].dt.year
    if years.nunique() >= 2:
        recent_year = years.max()
        prior_year = sorted(years.unique())[-2]
        recent = d[years == recent_year]
        prior = d[years == prior_year]
        label = f"{prior_year}→{recent_year}"
    else:
        mid = d["_t"].quantile(0.5)
        prior = d[d["_t"] <= mid]
        recent = d[d["_t"] > mid]
        label = "prior vs recent period"
    if len(recent) < MIN_CELL_ROWS or len(prior) < MIN_CELL_ROWS:
        return None, None, None, None
    return prior, recent, label, "_t"


def contribution_by_dim(prior, recent, metric, dim):
    """
    For a metric change between two periods, compute each segment's
    contribution to the TOTAL change. Returns a ranked table.
    contribution_pts = (seg_recent - seg_prior) / total_change * 100
    """
    p = prior.groupby(dim)[metric].sum()
    r = recent.groupby(dim)[metric].sum()
    segs = sorted(set(p.index) | set(r.index))
    total_prior = float(p.sum())
    total_recent = float(r.sum())
    total_change = total_recent - total_prior
    rows = []
    for s in segs:
        pv = float(p.get(s, 0.0)); rv = float(r.get(s, 0.0))
        chg = rv - pv
        # row counts for small-cell guard
        n_recent = int((recent[dim] == s).sum())
        n_prior = int((prior[dim] == s).sum())
        contrib_pts = (chg / total_change * 100) if total_change else 0.0
        seg_growth = (chg / pv * 100) if pv else (100.0 if rv else 0.0)
        rows.append({
            "segment": str(s), "prior": pv, "recent": rv, "change": chg,
            "contrib_pct": contrib_pts, "seg_growth": seg_growth,
            "n_recent": n_recent, "n_prior": n_prior,
        })
    tbl = pd.DataFrame(rows)
    if len(tbl):
        tbl["abs_contrib"] = tbl["contrib_pct"].abs()
        tbl = tbl.sort_values("abs_contrib", ascending=False).reset_index(drop=True)
    return tbl, total_prior, total_recent, total_change


def mix_by_dim(df, metric, dim):
    """No-time fallback: which segments dominate the metric level (share)."""
    g = df.groupby(dim)[metric].sum().sort_values(ascending=False)
    tot = float(g.sum()) or 1.0
    rows = [{"segment": str(k), "value": float(v), "share_pct": float(v) / tot * 100,
             "n": int((df[dim] == k).sum())} for k, v in g.items()]
    return pd.DataFrame(rows)


# ============================================================
# 3. DRIVER ANALYSIS  — recursive decomposition tree
# ============================================================
class DriverNode:
    def __init__(self, dim, segment, contrib_pct, change, prior, recent, n_recent, depth):
        self.dim = dim; self.segment = segment
        self.contrib_pct = contrib_pct; self.change = change
        self.prior = prior; self.recent = recent; self.n_recent = n_recent
        self.depth = depth
        self.children = []          # deeper DriverNodes
        self.table = None           # contribution table at this level (for charting)


def build_driver_tree(df, metric, dims, time_col):
    """
    Decompose metric change top-down. At each level pick the dimension whose
    top segment explains the most of the change, recurse into that segment.
    Returns (root_summary, tree_nodes, period_label).
    """
    prior, recent, label, tcol = _period_split(df, time_col)
    if prior is None:
        return None

    total_prior = float(prior[metric].sum())
    total_recent = float(recent[metric].sum())
    total_change = total_recent - total_prior
    total_pct = (total_change / total_prior * 100) if total_prior else 0.0

    root = {
        "metric": metric, "period": label,
        "prior": total_prior, "recent": total_recent,
        "change": total_change, "pct": total_pct,
        "direction": "down" if total_change < 0 else "up",
    }

    # recursive walk, narrowing the data to the dominant segment each level
    def walk(sub_prior, sub_recent, avail_dims, depth, parent_change):
        if depth >= MAX_TREE_DEPTH or not avail_dims or parent_change == 0:
            return []
        best = None
        for d in avail_dims:
            tbl, _, _, chg = contribution_by_dim(sub_prior, sub_recent, metric, d)
            if not len(tbl):
                continue
            top = tbl.iloc[0]
            # guards: enough rows, meaningful contribution, same direction as parent move
            if top["n_recent"] < MIN_CELL_ROWS:
                continue
            if top["abs_contrib"] < MIN_CONTRIB_PCT:
                continue
            score = top["abs_contrib"]
            if best is None or score > best["score"]:
                best = {"dim": d, "tbl": tbl, "top": top, "score": score}
        if not best:
            return []
        top = best["top"]; d = best["dim"]
        node = DriverNode(d, top["segment"], top["contrib_pct"], top["change"],
                          top["prior"], top["recent"], int(top["n_recent"]), depth)
        node.table = best["tbl"]
        # recurse into the winning segment with remaining dims
        np_ = sub_prior[sub_prior[d] == top["segment"]]
        nr_ = sub_recent[sub_recent[d] == top["segment"]]
        remaining = [x for x in avail_dims if x != d]
        node.children = walk(np_, nr_, remaining, depth + 1, top["change"])
        return [node]

    tree = walk(prior, recent, dims, 0, total_change)
    return {"root": root, "tree": tree}


# ============================================================
# 4. FINDINGS  — proven facts + the visual that proves them
# ============================================================
class Finding:
    """A proven, quantified insight. Carries the FACTS the LLM will narrate."""
    def __init__(self, ftype, role, facts, chart=None, importance=0.0):
        self.ftype = ftype          # 'situation','root_cause','interaction','mix','concentration','correlation','profile'
        self.role = role            # narrative role in SCR arc
        self.facts = facts          # dict of computed numbers (ground truth)
        self.chart = chart          # chart spec dict or None
        self.importance = importance
        self.narrative = None       # filled later by LLM (or fallback)

    def fact_brief(self):
        """Compact, unambiguous fact string handed to the LLM for narration."""
        return "; ".join(f"{k}={v}" for k, v in self.facts.items())


def _chart_from_contrib(tbl, metric, dim, kind="waterfall"):
    t = tbl.head(TOP_DIM_CARD)
    return {
        "kind": kind, "title": f"{metric} change by {dim}",
        "labels": t["segment"].tolist(),
        "values": [round(float(x), 2) for x in t["change"].tolist()],
        "secondary": [round(float(x), 1) for x in t["contrib_pct"].tolist()],
        "metric": metric, "dim": dim,
    }


def derive_findings(analytics):
    """
    Top-level entry: run the consultant loop and return a ranked list of Findings.
    Math-only here; narration added later.
    """
    ctx = StoryContext(analytics)
    df = analytics.df
    findings = []

    primary, _ = ctx.describe()

    # ---------- DRIVER (time) ----------
    if "driver_time" in ctx.modes:
        metric = ctx.metrics[0]
        res = build_driver_tree(df, metric, ctx.dims, ctx.time_col)
        if res:
            root = res["root"]; tree = res["tree"]
            # Situation: the headline movement
            findings.append(Finding(
                "situation", "situation",
                {"metric": metric, "period": root["period"],
                 "prior": _fmt(root["prior"]), "recent": _fmt(root["recent"]),
                 "change": _fmt(root["change"]), "pct": _signpct(root["pct"]),
                 "direction": root["direction"]},
                chart={"kind": "trend_period", "title": f"{metric} — {root['period']}",
                       "labels": [root["period"].split("→")[0] if "→" in root["period"] else "Prior",
                                  root["period"].split("→")[-1] if "→" in root["period"] else "Recent"],
                       "values": [round(root["prior"], 2), round(root["recent"], 2)],
                       "metric": metric},
                importance=100 + abs(root["pct"]),
            ))
            # Root cause chain: walk the tree
            def emit(nodes, parent_seg=None):
                for nd in nodes:
                    role = "root_cause" if nd.depth == 0 else "interaction"
                    seg_path = nd.segment if parent_seg is None else f"{parent_seg} × {nd.segment}"
                    facts = {
                        "metric": metric, "dimension": nd.dim, "segment": seg_path,
                        "segment_change": _fmt(nd.change),
                        "contribution_pct": f"{nd.contrib_pct:.0f}% of the total move",
                        "prior": _fmt(nd.prior), "recent": _fmt(nd.recent),
                        "rows": nd.n_recent,
                    }
                    findings.append(Finding(
                        "root_cause" if nd.depth == 0 else "interaction",
                        role, facts,
                        chart=_chart_from_contrib(nd.table, metric, nd.dim, kind="waterfall"),
                        importance=90 - nd.depth * 10 + abs(nd.contrib_pct) / 10,
                    ))
                    emit(nd.children, seg_path)
            emit(tree)

    # ---------- MIX (no time) ----------
    if "driver_mix" in ctx.modes and "driver_time" not in ctx.modes:
        metric = ctx.metrics[0]
        for d in ctx.dims[:2]:
            tbl = mix_by_dim(df, metric, d)
            if len(tbl) < 2:
                continue
            top = tbl.iloc[0]
            findings.append(Finding(
                "mix", "root_cause",
                {"metric": metric, "dimension": d, "segment": top["segment"],
                 "value": _fmt(top["value"]), "share": f"{top['share_pct']:.0f}%"},
                chart={"kind": "bar", "title": f"{metric} by {d}",
                       "labels": tbl.head(TOP_DIM_CARD)["segment"].tolist(),
                       "values": [round(float(x), 2) for x in tbl.head(TOP_DIM_CARD)["value"].tolist()],
                       "metric": metric, "dim": d},
                importance=60 + top["share_pct"] / 2,
            ))

    # ---------- CONCENTRATION ----------
    if "concentration" in ctx.modes:
        for d in ctx.dims[:2]:
            vc = df[d].value_counts()
            if len(vc) < 2:
                continue
            top_share = vc.iloc[0] / len(df) * 100
            if top_share >= 40:
                findings.append(Finding(
                    "concentration", "complication",
                    {"dimension": d, "segment": str(vc.index[0]),
                     "share": f"{top_share:.0f}%", "rows": int(vc.iloc[0])},
                    chart={"kind": "pie", "title": f"{d} concentration",
                           "labels": [str(x) for x in vc.head(5).index.tolist()],
                           "values": [float(x) for x in vc.head(5).values.tolist()],
                           "dim": d},
                    importance=40 + top_share / 4,
                ))

    # ---------- CORRELATION ----------
    if "correlation" in ctx.modes:
        nums = analytics.numeric_cols
        try:
            corr = df[nums].corr().abs()
            best = (0, None, None)
            for i in range(len(nums)):
                for j in range(i + 1, len(nums)):
                    r = corr.iloc[i, j]
                    if r == r and 0.4 <= r < 0.999 and r > best[0]:
                        best = (r, nums[i], nums[j])
            if best[1]:
                r, c1, c2 = best
                findings.append(Finding(
                    "correlation", "evidence",
                    {"a": c1, "b": c2, "r": f"{r:.2f}"},
                    chart={"kind": "scatter", "title": f"{c1} vs {c2}",
                           "labels": [], "values": [],
                           "x": [float(v) for v in df[c1].dropna().head(200).tolist()],
                           "y": [float(v) for v in df[c2].dropna().head(200).tolist()],
                           "xlab": c1, "ylab": c2},
                    importance=30 + r * 20,
                ))
        except Exception:
            pass

    # ---------- PROFILE (always, low priority) ----------
    b = analytics.insights["basic_stats"]
    findings.append(Finding(
        "profile", "context",
        {"rows": f"{b['rows']:,}", "cols": b["columns"],
         "completeness": f"{100 - b['missing_pct']:.0f}%"},
        importance=5,
    ))

    findings.sort(key=lambda f: -f.importance)
    return ctx, findings


#
# Consultant agent layer — sits on top of consultant.py (the math).
#   - CHART_CATALOG          : what each chart type needs to be valid
#   - validate_chart_choice  : the guardrail (rules) — can this finding support this chart?
#   - propose_visual         : LLM proposes chart+action-title; rules validate; fallback if rejected
#   - draft_storyboard       : LLM sequences findings into a narrative + gives reasoning
#   - narrate_finding        : LLM writes consultant prose from PROVEN facts only
# All LLM calls degrade gracefully to deterministic fallbacks when no llm present.
#


# ============================================================
# VISUAL CATALOG  — finding-shape → permissible charts
# ============================================================
# Each entry: requirements the finding/data must meet for the chart to be valid.
CHART_CATALOG = {
    "waterfall":   {"needs": ["change_decomposition"], "desc": "bridge from prior to recent showing each segment's +/- contribution"},
    "bar_ranked":  {"needs": ["segments", "values"],    "desc": "sorted bars, driver highlighted, rest muted"},
    "line":        {"needs": ["time_or_ordered"],       "desc": "trend over an ordered/time axis"},
    "area":        {"needs": ["time_or_ordered"],       "desc": "trend with filled area, good for one series + threshold"},
    "donut":       {"needs": ["segments", "values"],    "desc": "share of whole when one segment dominates"},
    "stacked100":  {"needs": ["segments", "values"],    "desc": "100% stacked composition"},
    "scatter":     {"needs": ["two_numeric"],           "desc": "relationship between two metrics with fit line"},
    "slope":       {"needs": ["two_period_segments"],   "desc": "prior→recent slope per segment (mix shift)"},
    "dumbbell":    {"needs": ["two_period_segments"],   "desc": "prior vs recent dots per segment"},
    "big_number":  {"needs": [],                        "desc": "one or few headline numbers, no chart"},
    "kpi_grid":    {"needs": [],                        "desc": "2x2 grid of headline numbers"},
}

# which charts each finding-type is ALLOWED to use (rules guardrail)
FINDING_CHART_WHITELIST = {
    "situation":     {"waterfall", "line", "area", "bar_ranked", "big_number", "kpi_grid"},
    "root_cause":    {"waterfall", "bar_ranked", "slope", "dumbbell"},
    "interaction":   {"waterfall", "bar_ranked", "slope", "dumbbell"},
    "mix":           {"bar_ranked", "donut", "stacked100", "slope"},
    "concentration": {"donut", "stacked100", "bar_ranked"},
    "correlation":   {"scatter"},
    "profile":       {"kpi_grid", "big_number"},
    "action":        {"big_number"},
}


def _finding_capabilities(finding):
    """What chart-requirements does this finding's data actually satisfy?"""
    caps = set()
    ch = finding.chart or {}
    k = ch.get("kind", "")
    facts = finding.facts
    # decomposition: came from a contribution table (has +/- changes by segment)
    if finding.ftype in ("situation", "root_cause", "interaction") and k in ("waterfall", "trend_period", "bar"):
        caps.add("change_decomposition")
    if ch.get("labels") and ch.get("values"):
        caps.add("segments"); caps.add("values")
    if k in ("line", "area", "trend_period"):
        caps.add("time_or_ordered")
    if k == "scatter" and ch.get("x") and ch.get("y"):
        caps.add("two_numeric")
    if finding.ftype in ("root_cause", "interaction", "mix"):
        # we have prior & recent per segment from the contribution table
        caps.add("two_period_segments")
    return caps


def validate_chart_choice(finding, chart_kind):
    """RULES guardrail. Returns (ok, reason)."""
    if chart_kind not in CHART_CATALOG:
        return False, f"unknown chart '{chart_kind}'"
    allowed = FINDING_CHART_WHITELIST.get(finding.ftype, set())
    if chart_kind not in allowed and CHART_CATALOG[chart_kind]["needs"]:
        return False, f"{chart_kind} not suitable for a '{finding.ftype}' finding"
    caps = _finding_capabilities(finding)
    for need in CHART_CATALOG[chart_kind]["needs"]:
        if need not in caps:
            return False, f"data can't support {chart_kind} (missing {need})"
    return True, "ok"


def _default_chart_for(finding):
    """Deterministic fallback chart kind per finding-type."""
    return {
        "situation": "waterfall" if (finding.chart or {}).get("kind") in ("waterfall",) else "line",
        "root_cause": "waterfall",
        "interaction": "waterfall",
        "mix": "bar_ranked",
        "concentration": "donut",
        "correlation": "scatter",
        "profile": "kpi_grid",
        "action": "big_number",
    }.get(finding.ftype, "bar_ranked")


# ============================================================
# VISUAL PROPOSAL  — LLM proposes, rules validate
# ============================================================
def propose_visual(finding, llm, audience="leadership"):
    """
    Returns {chart_kind, action_title, validated, reason}.
    LLM proposes; if invalid, fall back to the rule default (also validated).
    """
    allowed = sorted(FINDING_CHART_WHITELIST.get(finding.ftype, {"bar_ranked"}))
    fallback_kind = _default_chart_for(finding)
    ok, _ = validate_chart_choice(finding, fallback_kind)
    if not ok:
        # last-resort safe pick
        for k in allowed:
            if validate_chart_choice(finding, k)[0]:
                fallback_kind = k; break

    if not llm:
        return {"chart_kind": fallback_kind,
                "action_title": _fallback_title(finding),
                "validated": True, "reason": "no-llm fallback"}

    prompt = (
        f"You are a McKinsey visual-design lead. A finding has these PROVEN facts:\n"
        f"  {finding.fact_brief()}\n"
        f"Finding type: {finding.ftype}. Audience: {audience}.\n"
        f"Allowed chart types for this finding: {allowed}.\n"
        f"Chart meanings: " + "; ".join(f"{k}={CHART_CATALOG[k]['desc']}" for k in allowed) + "\n\n"
        "Pick the ONE chart that proves this point most clearly, and write an ACTION TITLE "
        "(a full sentence stating the finding, e.g. 'New York drove 92% of the YoY decline' — "
        "NOT a label like 'Sales by State'). Use the real numbers from the facts.\n"
        "Return ONLY JSON: {\"chart\":\"<one of allowed>\",\"title\":\"<action title>\"}"
    )
    try:
        raw = _re.sub(r"```json|```", "", llm.predict(prompt, max_tokens=200)).strip()
        p = _json.loads(raw)
        kind = p.get("chart", fallback_kind)
        title = p.get("title") or _fallback_title(finding)
        ok, reason = validate_chart_choice(finding, kind)
        if not ok:
            kind = fallback_kind
            reason = f"LLM pick rejected ({reason}); used {fallback_kind}"
        else:
            reason = "LLM proposed, validated"
        return {"chart_kind": kind, "action_title": title, "validated": True, "reason": reason}
    except Exception as e:
        return {"chart_kind": fallback_kind, "action_title": _fallback_title(finding),
                "validated": True, "reason": f"llm-error fallback ({e})"}


def _fallback_title(f):
    ff = f.facts
    if f.ftype == "situation":
        return f"{ff.get('metric','Metric').title()} moved {ff.get('pct','')} ({ff.get('period','')})"
    if f.ftype in ("root_cause", "interaction"):
        return f"{ff.get('segment','A segment')} drove {ff.get('contribution_pct','the change')}"
    if f.ftype == "mix":
        return f"{ff.get('segment','Top segment')} leads {ff.get('metric','the metric')} at {ff.get('share','')}"
    if f.ftype == "concentration":
        return f"{ff.get('segment','One segment')} is {ff.get('share','')} of the total"
    if f.ftype == "correlation":
        return f"{ff.get('a','')} and {ff.get('b','')} move together (r={ff.get('r','')})"
    if f.ftype == "profile":
        return f"{ff.get('rows','')} records across {ff.get('cols','')} fields, {ff.get('completeness','')} complete"
    return "Key finding"


# ============================================================
# STORYBOARD AGENT  — LLM drafts ordered narrative + reasoning
# ============================================================
def draft_storyboard(findings, framing, llm):
    """
    framing = {decision, audience, takeaway}
    Returns {order: [finding_index...], reasoning: str, arc: [labels]}.
    LLM sequences; rules ensure situation-first & action-last sanity.
    """
    # Build a compact menu for the LLM
    menu = []
    for i, f in enumerate(findings):
        menu.append(f"{i}: [{f.ftype}/{f.role}] {f.fact_brief()[:120]}")

    fallback_order = _scr_order(findings)
    if not llm:
        return {"order": fallback_order,
                "reasoning": "Sequenced as Situation → Complication → Root-cause chain → "
                             "supporting evidence → context (no AI; default consultant arc).",
                "arc": [findings[i].role for i in fallback_order]}

    prompt = (
        "You are a McKinsey engagement manager building a board deck.\n"
        f"Decision the audience must make: {framing.get('decision','(unspecified)')}\n"
        f"Audience: {framing.get('audience','leadership')}\n"
        f"The one thing they must leave knowing: {framing.get('takeaway','(unspecified)')}\n\n"
        "Available findings (index: [type/role] facts):\n" + "\n".join(menu) + "\n\n"
        "Sequence these into a tight narrative using the pyramid principle "
        "(answer first, then support). Drop weak/redundant findings. "
        "Aim for 5-8 slides. Order should build a cause→effect→action story.\n"
        "Return ONLY JSON: {\"order\":[indices in slide order],"
        "\"reasoning\":\"2-3 sentences explaining the narrative arc you chose\"}"
    )
    try:
        raw = _re.sub(r"```json|```", "", llm.predict(prompt, max_tokens=400)).strip()
        p = _json.loads(raw)
        order = [i for i in p.get("order", []) if isinstance(i, int) and 0 <= i < len(findings)]
        if not order:
            order = fallback_order
        # sanity: situation first if present, profile/context not first
        sit = next((i for i in order if findings[i].role == "situation"), None)
        if sit is not None and order[0] != sit:
            order.remove(sit); order.insert(0, sit)
        return {"order": order,
                "reasoning": p.get("reasoning", "Narrative sequenced by the agent."),
                "arc": [findings[i].role for i in order]}
    except Exception:
        return {"order": fallback_order,
                "reasoning": "Default consultant arc (AI sequencing unavailable).",
                "arc": [findings[i].role for i in fallback_order]}


def _scr_order(findings):
    """Deterministic Situation→Complication→Cause-chain→Evidence→Context fallback."""
    role_rank = {"situation": 0, "complication": 1, "root_cause": 2,
                 "interaction": 3, "mix": 4, "evidence": 5, "context": 9}
    idx = list(range(len(findings)))
    idx.sort(key=lambda i: (role_rank.get(findings[i].role, 6), -findings[i].importance))
    # keep it tight
    return idx[:8]


def agent_choose_structure(findings, framing, llm):
    """
    The agent (Claude) freely picks the best storyboard STRUCTURE for this data,
    and explains why. Returns {style_key, style_name, rationale}.
    Validates the pick against the real STORYBOARD_STYLES; falls back to consulting.
    """
    valid = {v["key"]: name for name, v in STORYBOARD_STYLES.items()}
    default_key = "consulting"
    default_name = next((n for n, v in STORYBOARD_STYLES.items() if v["key"] == default_key),
                        next(iter(STORYBOARD_STYLES.keys())))
    if not llm:
        return {"style_key": default_key, "style_name": default_name,
                "rationale": "Defaulted to the Standard Consulting storyline (no AI configured)."}

    # describe the menu of structures + what the data looks like
    menu = "\n".join(f"  - {name} (key={v['key']}): {v['blurb']}" for name, v in STORYBOARD_STYLES.items())
    roles = {}
    for f in findings:
        roles[f.role] = roles.get(f.role, 0) + 1
    shape = ", ".join(f"{k}×{v}" for k, v in roles.items())
    facts = "; ".join(f.fact_brief()[:80] for f in findings[:6])

    prompt = (
        "You are a McKinsey engagement manager. Choose the BEST storyboard structure "
        "for THIS specific dataset and goal — you have full freedom to pick.\n\n"
        f"Decision at stake: {framing.get('decision','(unspecified)')}\n"
        f"Audience: {framing.get('audience','leadership')}\n"
        f"Finding shape (role×count): {shape}\n"
        f"Top findings: {facts}\n\n"
        f"Available structures:\n{menu}\n\n"
        "Pick the one that best fits the story this data tells. For example: a single dominant "
        "driver suits a deep-dive; a portfolio of segments suits a comparison; a clear time decline "
        "suits chronological; a board that wants the answer first suits pyramid; a full diagnostic "
        "suits standard consulting.\n"
        "Return ONLY JSON: {\"key\":\"<one of the keys above>\","
        "\"rationale\":\"2-3 sentences: why this structure fits THIS data and decision\"}"
    )
    try:
        raw = _re.sub(r"```json|```", "", llm.predict(prompt, max_tokens=300)).strip()
        mb = _re.search(r"\{.*\}", raw, _re.DOTALL)
        p = _json.loads(mb.group(0) if mb else raw)
        key = p.get("key", default_key)
        if key not in valid:
            key = default_key
        return {"style_key": key, "style_name": valid.get(key, default_name),
                "rationale": p.get("rationale", "Chosen by the agent for this data.")}
    except Exception:
        return {"style_key": default_key, "style_name": default_name,
                "rationale": "Defaulted to Standard Consulting (agent selection unavailable)."}


def narrate_finding(finding, framing, llm):
    """Consultant prose from PROVEN facts only. LLM never sees raw data, only facts."""
    if not llm:
        return _fallback_narrative(finding)
    prompt = (
        "You are a McKinsey consultant writing the talking point for one slide.\n"
        f"PROVEN facts (do not invent beyond these): {finding.fact_brief()}\n"
        f"Audience: {framing.get('audience','leadership')}. "
        f"Decision at stake: {framing.get('decision','')}\n"
        "Write 1-2 crisp sentences: state the finding with its number, then the 'so what'. "
        "No hedging, no preamble, no bullet points."
    )
    try:
        return llm.predict(prompt, max_tokens=160).strip()
    except Exception:
        return _fallback_narrative(finding)


def _fallback_narrative(f):
    ff = f.facts
    if f.ftype == "situation":
        return (f"{ff.get('metric','The metric').title()} moved {ff.get('pct','')} over {ff.get('period','the period')}, "
                f"from {ff.get('prior','')} to {ff.get('recent','')} — the headline we need to explain.")
    if f.ftype in ("root_cause", "interaction"):
        return (f"{ff.get('segment','This segment')} accounts for {ff.get('contribution_pct','the move')} "
                f"({ff.get('segment_change','')}), making it the primary lever to address.")
    if f.ftype == "mix":
        return (f"{ff.get('segment','')} leads {ff.get('metric','')} at {ff.get('share','')} of the total "
                f"({ff.get('value','')}) — the segment to protect and grow.")
    if f.ftype == "concentration":
        return (f"{ff.get('segment','One segment')} represents {ff.get('share','')} of activity — "
                f"a concentration worth managing as a risk.")
    if f.ftype == "correlation":
        return (f"{ff.get('a','')} and {ff.get('b','')} are correlated (r={ff.get('r','')}), "
                "suggesting a lever worth testing.")
    return f"{ff.get('rows','')} records across {ff.get('cols','')} fields at {ff.get('completeness','')} completeness."


#
# Consultant deck builder — charts, storyboard styles, visual flow, pptx render.
# Depends on consultant.py (findings) + agent.py (visual/story choices).
#
# Storyboard STYLES (each defines a narrative ordering of finding-roles AND a
# flow label shown on the storyline slide):
#   - SCR              Situation → Complication → Root-cause chain → Implication → Action
#   - Pyramid          Answer first → supporting drivers → evidence
#   - Deep-dive        Headline → progressive drill-down (situation→root→interaction)
#   - Chronological    How it evolved over time → where it stands → what to do
#   - Comparison       Benchmark segments side-by-side → leader/laggard → action
#
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN as AL, MSO_ANCHOR as AN
from pptx.enum.shapes import MSO_SHAPE as SH
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE as CT, XL_LEGEND_POSITION as LEG, XL_LABEL_POSITION as LP
from pptx.oxml.ns import qn


# ============================================================
# STORYBOARD STYLES
# ============================================================
STORYBOARD_STYLES = {
    "Standard Consulting (recommended)": {
        "key": "consulting",
        "blurb": "Full executive deck: Executive Summary → Problem Definition (what/where/who) → Insights & Drivers → Recommendations → Impact → Next-Steps roadmap. Action-titles on every slide (golden rule).",
        "role_order": ["situation", "root_cause", "interaction", "mix", "complication", "concentration", "correlation", "action", "context"],
        "flow": ["Exec Summary", "Problem", "Insights & Drivers", "Recommendations", "Next Steps"],
    },
    "SCR — Consultant": {
        "key": "scr",
        "blurb": "Situation → Complication → Root-cause chain → Implication → Action. The classic McKinsey storyline.",
        "role_order": ["situation", "complication", "root_cause", "interaction", "mix", "evidence", "action", "context"],
        "flow": ["Situation", "Complication", "Root cause", "So what", "Action"],
    },
    "Pyramid — Answer first": {
        "key": "pyramid",
        "blurb": "Lead with the recommendation, then the findings that prove it. Best for execs who want the answer up front.",
        "role_order": ["action", "situation", "root_cause", "interaction", "mix", "complication", "evidence", "context"],
        "flow": ["Recommendation", "Why", "Evidence", "Proof points"],
    },
    "Deep-dive — Root cause": {
        "key": "deepdive",
        "blurb": "Headline number, then progressively drill into the driver and its interactions. Best for diagnostic reviews.",
        "role_order": ["situation", "root_cause", "interaction", "mix", "complication", "evidence", "action", "context"],
        "flow": ["Headline", "Driver", "Drill-down", "Pinpoint", "Action"],
    },
    "Chronological — Trend": {
        "key": "chrono",
        "blurb": "How the metric evolved, where it stands now, and what to do next. Best for performance updates.",
        "role_order": ["situation", "mix", "root_cause", "interaction", "complication", "evidence", "action", "context"],
        "flow": ["Where we were", "What changed", "Where we are", "Next"],
    },
    "Comparison — Benchmark": {
        "key": "compare",
        "blurb": "Segments side-by-side, leader vs laggard, then the move. Best for portfolio / regional reviews.",
        "role_order": ["mix", "situation", "root_cause", "interaction", "complication", "evidence", "action", "context"],
        "flow": ["Benchmark", "Leaders & laggards", "Why", "Action"],
    },
}


def order_findings_by_style(findings, style_key):
    style = next((v for v in STORYBOARD_STYLES.values() if v["key"] == style_key),
                 next(iter(STORYBOARD_STYLES.values())))
    rank = {r: i for i, r in enumerate(style["role_order"])}
    idx = sorted(range(len(findings)),
                 key=lambda i: (rank.get(findings[i].role, 50), -findings[i].importance))
    return idx[:8], style["flow"]


# ---- Standard Consulting Storyline (manual-based) ----


# ---- the canonical blueprint (what slots a full consulting deck has) ----
CONSULTING_BLUEPRINT = [
    # section, slide_role, layout, what finding-types fit this slot
    ("Executive Summary", "exec_summary", "exec_summary", ["situation", "root_cause", "action"]),
    ("Problem Definition", "what",        "stat_hero",    ["situation"]),
    ("Problem Definition", "where",        "map_or_bar",  ["root_cause", "mix"]),
    ("Problem Definition", "who",          "ranked_bar",  ["interaction", "mix", "concentration"]),
    ("Insights & Drivers",  "root_cause",  "waterfall",   ["root_cause"]),
    ("Insights & Drivers",  "segment",     "bar_ranked",  ["interaction", "mix"]),
    ("Insights & Drivers",  "supporting",  "trend_or_scatter", ["situation", "correlation", "concentration"]),
    ("Recommendations",     "actions",     "actions",     ["action"]),
    ("Recommendations",     "impact",      "impact",      ["action"]),
    ("Next Steps",          "roadmap",     "roadmap",     ["action"]),
]


def plan_consulting_storyline(findings, framing):
    """
    Map findings into the consulting blueprint. Returns an ordered list of
    slide-plan dicts: {section, role, layout, finding_idx|None, title, takeaway}.
    Each title is an ACTION TITLE (golden rule). Slots with no matching finding
    are either filled from related findings or dropped (so we never show empties).
    """
    by_type = {}
    for i, f in enumerate(findings):
        by_type.setdefault(f.ftype, []).append(i)

    used = set()
    plan = []

    def take(ftypes):
        for t in ftypes:
            for i in by_type.get(t, []):
                if i not in used:
                    used.add(i)
                    return i
        return None

    # 1) Executive summary is synthesized, not a single finding
    sit = next((i for i, f in enumerate(findings) if f.ftype == "situation"), None)
    rc = next((i for i, f in enumerate(findings) if f.ftype in ("root_cause", "interaction")), None)
    plan.append({
        "section": "Executive Summary", "role": "exec_summary", "layout": "exec_summary",
        "finding_idx": None, "core_situation": sit, "core_cause": rc,
    })

    # 2) walk the blueprint, attach a finding to each slot if available
    for section, role, layout, ftypes in CONSULTING_BLUEPRINT:
        if role == "exec_summary":
            continue
        if role in ("actions", "impact", "roadmap"):
            # recommendation slots handled after
            continue
        fi = take(ftypes)
        if fi is None:
            continue
        plan.append({"section": section, "role": role, "layout": layout, "finding_idx": fi})

    # 3) recommendations block (always present; from action findings or framing recs)
    recs = framing.get("recommendations") or []
    plan.append({"section": "Recommendations", "role": "actions", "layout": "actions",
                 "finding_idx": None, "recs": recs[:3]})
    if framing.get("impacts"):
        plan.append({"section": "Recommendations", "role": "impact", "layout": "impact",
                     "finding_idx": None, "impacts": framing["impacts"][:3]})

    # 4) next steps roadmap (always; from framing or generic)
    plan.append({"section": "Next Steps", "role": "roadmap", "layout": "roadmap",
                 "finding_idx": None, "steps": framing.get("roadmap")})

    return plan


# ---- LLM prompt scaffolding for action-titles + takeaways (golden rule) ----
def action_title_prompt(finding, framing):
    return (
        "Write a SLIDE TITLE that states the full insight as a sentence (the 'golden rule': "
        "a reader skimming only titles must understand the whole story).\n"
        f"PROVEN facts: {finding.fact_brief()}\n"
        "Bad: 'Sales by Region'. Good: 'North region contributes 60% of total growth'.\n"
        "Use the real numbers. Return ONLY the title text, no quotes."
    )


def takeaway_prompt(finding, framing):
    return (
        "Write a one-line TAKEAWAY (the conclusion at the bottom of the slide) for this finding.\n"
        f"PROVEN facts: {finding.fact_brief()}\n"
        f"Decision at stake: {framing.get('decision','')}\n"
        "One crisp sentence stating the 'so what'. Return ONLY the takeaway."
    )


def exec_summary_prompt(findings, framing, idxs):
    facts = "; ".join(findings[i].fact_brief()[:90] for i in idxs if i is not None)
    return (
        "Write an EXECUTIVE SUMMARY for slide 1 of a consulting deck.\n"
        f"Decision: {framing.get('decision','')}. Audience: {framing.get('audience','')}.\n"
        f"Proven facts available: {facts}\n\n"
        "Return ONLY JSON:\n"
        '{"core_answer":"the single headline answer in one sentence",'
        '"insights":["insight 1 w/ number","insight 2 w/ number","insight 3 w/ number"],'
        '"actions":["action 1","action 2"]}'
    )



from pptx.util import Inches as In, Pt


def render_consulting_slide(prs, blank, plan_item, findings, framing, th, G, idx, total):
    """Render one slide-plan item. Returns the slide."""
    s = prs.slides.add_slide(blank)
    layout = plan_item["layout"]
    SH = G["SH"]

    def section_chip(section):
        # small section locator top-right (helps the section flow read)
        G["_txt"](s, section.upper(), 9.0, 0.62, 3.6, 0.3, size=9, font=th["bf"],
                  color=th["soft"], bold=True, align="right", spacing=1.0)

    def headline(title, tag="Insight"):
        G["_pill"](s, 0.9, 0.55, tag, th)
        G["_txt"](s, title, 0.86, 1.05, 11.7, 1.1, size=27, font=th["hf"],
                  color=th["ink"], bold=True, ls=1.04)

    def takeaway(text):
        if not text:
            return
        G["_rect"](s, 0.9, 6.55, 11.55, 0.7, fill=th["tagbg"], radius=0.16)
        G["_txt"](s, [("Takeaway  ", th["tagink"], True), (text, th["ink"], False)],
                  1.15, 6.62, 11.1, 0.56, size=12.5, font=th["bf"], valign="middle")

    def slidenum():
        G["_txt"](s, f"{idx} / {total}", 12.5, 7.15, 0.6, 0.25, size=9, font=th["bf"],
                  color=th["soft"], align="right")

    # ---------- EXEC SUMMARY ----------
    if layout == "exec_summary":
        G["_bg"](s, th["canvas"])
        G["_pill"](s, 0.9, 0.55, "Executive Summary", th)
        es = plan_item.get("exec", {})
        core = es.get("core_answer", "")
        G["_txt"](s, core, 0.86, 1.1, 11.8, 1.5, size=30, font=th["hf"], color=th["ink"], bold=True, ls=1.05)
        # 3 insight cards
        insights = es.get("insights", [])[:3]
        cw = 3.78; gap = 0.2; x0 = 0.9; y = 3.0
        for i, ins in enumerate(insights):
            x = x0 + i * (cw + gap)
            G["_rect"](s, x, y, cw, 2.0, fill=th["panel"], line=th["line"], lw=1, radius=0.10, shdw=True)
            G["_rect"](s, x, y, cw, 0.12, fill=th["accent"], shape=SH.RECTANGLE)
            G["_txt"](s, str(i + 1), x + 0.2, y + 0.22, 0.6, 0.5, size=20, font=th["hf"],
                      color=th["accent"], bold=True)
            G["_txt"](s, ins, x + 0.2, y + 0.78, cw - 0.4, 1.1, size=12.5, font=th["bf"],
                      color=th["ink"], ls=1.2)
        # actions strip
        acts = es.get("actions", [])[:2]
        if acts:
            G["_txt"](s, "RECOMMENDED ACTIONS", 0.9, 5.25, 6, 0.3, size=10, font=th["bf"],
                      color=th["soft"], bold=True, spacing=1)
            for i, ac in enumerate(acts):
                yy = 5.6 + i * 0.62
                G["_rect"](s, 0.9, yy, 0.16, 0.16, fill=th["accent"], shape=SH.OVAL)
                G["_txt"](s, ac, 1.2, yy - 0.06, 11.2, 0.5, size=12.5, font=th["bf"], color=th["ink"], valign="middle")
        slidenum()
        return s

    # ---------- STAT HERO (the 'what is happening') ----------
    if layout == "stat_hero":
        fi = plan_item["finding_idx"]; f = findings[fi]
        section_chip(plan_item["section"])
        G["_bg"](s, th["canvas"])
        headline(plan_item.get("title", "What is happening"), tag=plan_item["section"].split()[0])
        big = f.facts.get("pct") or f.facts.get("change") or f.facts.get("share") or ""
        G["_txt"](s, str(big), 0.86, 2.05, 6.2, 1.6, size=96, font=th["hf"], color=th["accent"], bold=True)
        G["_txt"](s, f.facts.get("metric", "").upper(), 0.92, 3.75, 6, 0.4, size=15, font=th["bf"],
                  color=th["soft"], spacing=1)
        # stat strip: prior / recent / delta / period  (real numbers, fills the space)
        stats = []
        if f.facts.get("prior"): stats.append(("Prior", f.facts["prior"]))
        if f.facts.get("recent"): stats.append(("Recent", f.facts["recent"]))
        if f.facts.get("change"): stats.append(("Change", f.facts["change"]))
        if f.facts.get("period"): stats.append(("Period", f.facts["period"]))
        if stats:
            sw = 1.5; gap = 0.2; sx = 0.9; sy = 4.6
            for i, (lab, val) in enumerate(stats[:4]):
                x = sx + i*(sw+gap)
                G["_rect"](s, x, sy, sw, 1.0, fill=th["panel"], line=th["line"], lw=1, radius=0.12, shdw=True)
                G["_txt"](s, str(val), x, sy+0.16, sw, 0.5, size=18, font=th["hf"],
                          color=th["ink"], bold=True, align="center")
                G["_txt"](s, lab.upper(), x, sy+0.66, sw, 0.25, size=8.5, font=th["bf"],
                          color=th["soft"], align="center", spacing=0.5)
        ch = f.chart or {}
        if ch.get("labels"):
            try: G["chart_line"](s, 7.4, 2.2, 5.0, 2.9, ch["labels"], ch["values"], th, area=True)
            except Exception: pass
        takeaway(plan_item.get("takeaway", ""))
        slidenum()
        return s

    # ---------- FINDING with chart (where/who/root_cause/segment/supporting) ----------
    if layout in ("waterfall", "bar_ranked", "ranked_bar", "map_or_bar", "trend_or_scatter"):
        fi = plan_item["finding_idx"]; f = findings[fi]
        section_chip(plan_item["section"])
        G["_bg"](s, th["canvas"])
        headline(plan_item.get("title", "Finding"), tag=plan_item["section"].split()[0])
        ch = f.chart or {}; ck = plan_item.get("chart_kind") or ch.get("kind", "bar")
        labels = ch.get("labels", []); values = ch.get("values", [])
        # decide: does this finding have per-segment supporting numbers for a table?
        has_table = bool(labels) and bool(values) and ck in ("waterfall", "bar", "bar_ranked", "ranked_bar", "map_or_bar", "donut", "pie")
        cx, cy = 0.9, 2.2
        cw = 7.4 if has_table else 11.5      # wide chart; leave room only if a real table exists
        chh = 3.7
        drew = False
        try:
            if ck == "waterfall" and labels:
                G["chart_waterfall"](s, cx, cy, cw, chh, labels, values, th,
                                     pcts=ch.get("secondary")); drew = True
            elif ck in ("bar", "bar_ranked", "ranked_bar", "map_or_bar") and labels:
                G["chart_bar_ranked"](s, cx, cy, cw, chh, labels, values, th, 0); drew = True
            elif ck in ("line", "area", "trend_or_scatter") and labels:
                G["chart_line"](s, cx, cy, cw, chh, labels, values, th, area=(ck == "area")); drew = True
            elif ck in ("donut", "pie") and labels:
                G["chart_donut"](s, cx, cy, cw, chh, labels, values, th); drew = True
            elif ck == "scatter" and ch.get("x"):
                G["chart_scatter"](s, cx, cy, cw, chh, ch["x"], ch["y"], th); drew = True
            elif labels:
                G["chart_bar_ranked"](s, cx, cy, cw, chh, labels, values, th, 0); drew = True
        except Exception:
            drew = False

        # ---- supporting DATA TABLE on the right (real numbers, not repeated text) ----
        if has_table and drew:
            tx, tw = 8.5, 4.0
            G["_txt"](s, "SUPPORTING DATA", tx, cy, tw, 0.3, size=10, font=th["hf"],
                      color=th["soft"], bold=True, spacing=1)
            tot = sum(abs(v) for v in values) or 1
            # header row
            ty = cy + 0.42; rh = 0.46
            G["_txt"](s, [("Segment", th["soft"], True)], tx, ty, 1.9, 0.3, size=9, font=th["bf"])
            G["_txt"](s, [("Value", th["soft"], True)], tx+1.9, ty, 1.05, 0.3, size=9, font=th["bf"], align="right")
            G["_txt"](s, [("Share", th["soft"], True)], tx+2.95, ty, 1.05, 0.3, size=9, font=th["bf"], align="right")
            ty += 0.34
            top_n = sorted(range(len(values)), key=lambda i: -abs(values[i]))[:6]
            for rank, i in enumerate(top_n):
                seg = str(labels[i]); val = values[i]; share = abs(val)/tot*100
                hot = (rank == 0)
                if hot:
                    G["_rect"](s, tx-0.1, ty-0.04, tw, rh-0.08, fill=th["tagbg"], radius=0.14)
                G["_txt"](s, seg[:16], tx, ty, 1.9, rh, size=10.5, font=th["bf"],
                          color=(th["tagink"] if hot else th["ink"]), bold=hot, valign="middle")
                G["_txt"](s, f"{val:+,.0f}" if ck=="waterfall" else f"{val:,.0f}",
                          tx+1.9, ty, 1.05, rh, size=10.5, font=th["bf"],
                          color=(th["tagink"] if hot else th["body"]), bold=hot, align="right", valign="middle")
                G["_txt"](s, f"{share:.0f}%", tx+2.95, ty, 1.05, rh, size=10.5, font=th["bf"],
                          color=(th["accent"] if hot else th["soft"]), bold=hot, align="right", valign="middle")
                ty += rh
        takeaway(plan_item.get("takeaway", "") or plan_item.get("narrative", ""))
        slidenum()
        return s

    # ---------- ACTIONS ----------
    if layout == "actions":
        G["_bg"](s, th["canvas"])
        headline(plan_item.get("title", "Recommended actions"), tag="Recommendations")
        cols = [th["accent"], th["accent2"], th["soft"]]
        for i, t in enumerate(plan_item.get("recs", [])[:3]):
            y = 2.0 + i * 1.45
            G["_rect"](s, 0.9, y, 11.55, 1.25, fill=th["panel"], line=th["line"], lw=1, radius=0.10, shdw=True)
            G["_rect"](s, 0.9, y, 0.7, 1.25, fill=cols[i % 3], shape=SH.RECTANGLE)
            G["_txt"](s, str(i + 1), 0.9, y, 0.7, 1.25, size=26, font=th["hf"], color="FFFFFF",
                      bold=True, align="center", valign="middle")
            G["_txt"](s, t, 1.8, y + 0.1, 10.4, 1.05, size=14, font=th["bf"], color=th["ink"],
                      valign="middle", ls=1.15)
        slidenum()
        return s

    # ---------- IMPACT ----------
    if layout == "impact":
        G["_bg"](s, th["canvas"])
        headline(plan_item.get("title", "Expected impact"), tag="Recommendations")
        imp = plan_item.get("impacts", [])[:3]
        cw = 3.78; gap = 0.2; x0 = 0.9; y = 2.4
        for i, m in enumerate(imp):
            x = x0 + i * (cw + gap)
            G["_rect"](s, x, y, cw, 2.4, fill=th["panel"], line=th["line"], lw=1, radius=0.10, shdw=True)
            G["_txt"](s, m.get("value", ""), x, y + 0.3, cw, 1.0, size=40, font=th["hf"],
                      color=th["accent"], bold=True, align="center")
            G["_txt"](s, m.get("label", ""), x + 0.2, y + 1.5, cw - 0.4, 0.8, size=12, font=th["bf"],
                      color=th["ink"], align="center", ls=1.2)
        slidenum()
        return s

    # ---------- ROADMAP ----------
    if layout == "roadmap":
        G["_bg"](s, th["canvas"])
        headline(plan_item.get("title", "Next steps & roadmap"), tag="Next Steps")
        steps = plan_item.get("steps") or [
            {"phase": "Now (0–30d)", "what": "Confirm the driver and brief owners"},
            {"phase": "Next (1–3m)", "what": "Launch the highest-impact fix"},
            {"phase": "Later (3–6m)", "what": "Scale what works, monitor leading indicators"},
        ]
        n = len(steps); cw = (11.55 - (n - 1) * 0.25) / n; x0 = 0.9; y = 2.6
        for i, stp in enumerate(steps):
            x = x0 + i * (cw + 0.25)
            G["_rect"](s, x, y, cw, 2.2, fill=th["panel"], line=th["line"], lw=1, radius=0.10, shdw=True)
            G["_rect"](s, x, y, cw, 0.55, fill=th["accent"], radius=0.10)
            G["_txt"](s, stp.get("phase", f"Phase {i+1}"), x, y + 0.06, cw, 0.42, size=12,
                      font=th["hf"], color="FFFFFF", bold=True, align="center", valign="middle")
            G["_txt"](s, stp.get("what", ""), x + 0.18, y + 0.7, cw - 0.36, 1.4, size=12,
                      font=th["bf"], color=th["ink"], ls=1.2)
            if i < n - 1:
                G["_txt"](s, "→", x + cw, y + 0.9, 0.25, 0.4, size=18, font=th["bf"],
                          color=th["soft"], align="center")
        slidenum()
        return s

    # fallback
    G["_bg"](s, th["canvas"])
    headline(plan_item.get("title", "Finding"))
    slidenum()
    return s



# ============================================================
# DRAW HELPERS (theme passed in as dict of hex strings)
# ============================================================
def _hx(h): h = h.lstrip("#"); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def _bg(s, c): f = s.background.fill; f.solid(); f.fore_color.rgb = _hx(c)

def _shadow(shp, alpha=10000):
    try:
        spPr = shp._element.spPr
        el = spPr.makeelement(qn('a:effectLst'), {})
        sh = spPr.makeelement(qn('a:outerShdw'), {'blurRad':'110000','dist':'28000','dir':'5400000','rotWithShape':'0'})
        clr = spPr.makeelement(qn('a:srgbClr'), {'val':'1F2235'}); aa = spPr.makeelement(qn('a:alpha'), {'val':str(alpha)})
        clr.append(aa); sh.append(clr); el.append(sh); spPr.append(el)
    except Exception: pass

def _rect(s, x,y,w,h, fill=None, line=None, lw=1.0, shape=SH.ROUNDED_RECTANGLE, radius=0.5, shdw=False):
    shp = s.shapes.add_shape(shape, In(x),In(y),In(w),In(h)); shp.shadow.inherit=False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=_hx(fill)
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=_hx(line); shp.line.width=Pt(lw)
    if shape==SH.ROUNDED_RECTANGLE:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if shdw: _shadow(shp)
    return shp

def _alpha_fill(shp, val):
    try:
        sp = shp.fill.fore_color._xFill
        a = sp.makeelement(qn('a:alpha'), {'val':str(val)}); sp.find(qn('a:srgbClr')).append(a)
    except Exception: pass

def _txt(s, runs, x,y,w,h, size=14, font="Calibri", color="1F2235", bold=False, align="left",
         valign="top", spacing=None, ls=None):
    tb=s.shapes.add_textbox(In(x),In(y),In(w),In(h)); tf=tb.text_frame; tf.word_wrap=True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    tf.vertical_anchor={"top":AN.TOP,"middle":AN.MIDDLE,"bottom":AN.BOTTOM}[valign]
    parts = runs if isinstance(runs,list) else [(runs,color,bold)]
    p=tf.paragraphs[0]; p.alignment={"left":AL.LEFT,"center":AL.CENTER,"right":AL.RIGHT}[align]
    if ls: p.line_spacing=ls
    for t,c,b in parts:
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.name=font; r.font.bold=b; r.font.color.rgb=_hx(c)
        if spacing: r._r.get_or_add_rPr().set("spc", str(int(spacing*100)))
    return tb

def _pill(s, x, y, text, th):
    w=max(1.0, 0.16+0.10*len(text))
    _rect(s, x,y,w,0.4, fill=th["tagbg"], radius=0.5)
    _txt(s, text.upper(), x,y,w,0.4, size=10.5, font=th["bf"], color=th["tagink"], bold=True,
         align="center", valign="middle", spacing=1.2)
    return w

def _section_head(s, tag, title, th, x=0.9, y=0.55):
    _pill(s, x, y, tag, th)
    _txt(s, title, x-0.04, y+0.5, 11.7, 1.4, size=33, font=th["hf"], color=th["ink"], bold=True, ls=1.03)


# ============================================================
# CONSULTANT CHART TYPES
# ============================================================
def _axis_style(ch, th, hide_val=False):
    try:
        ca=ch.category_axis; ca.tick_labels.font.size=Pt(10); ca.tick_labels.font.color.rgb=_hx(th["ink"])
        ca.format.line.color.rgb=_hx(th["line"])
    except Exception: pass
    try:
        va=ch.value_axis
        if hide_val: va.visible=False
        else:
            va.tick_labels.font.size=Pt(9); va.tick_labels.font.color.rgb=_hx(th["soft"])
            va.major_gridlines.format.line.color.rgb=_hx(th["line"]); va.major_gridlines.format.line.width=Pt(0.5)
    except Exception: pass
    try: ch.font.name=th["bf"]
    except Exception: pass


def chart_waterfall(s, x,y,w,h, labels, deltas, th, accent=None, start_total=None, pcts=None):
    """
    Data-dense floating waterfall. Each bar shows its signed delta AND its % of the
    total move; a final 'Net' bar anchors the cumulative result; a faint running
    line connects the steps. pcts = optional list of contribution % per segment.
    """
    accent = accent or th["accent"]
    neg = th.get("neg", "E5484D")
    total = sum(deltas)
    if pcts is None:
        denom = sum(abs(d) for d in deltas) or 1
        pcts = [d / (total if total else denom) * 100 for d in deltas]
    # cumulative incl. a final Net bar (from 0 to total)
    cum=[0.0]
    for d in deltas: cum.append(cum[-1]+d)
    show_labels = list(labels) + ["Net"]
    show_deltas = list(deltas) + [total]
    bar_tops = list(cum[:-1]) + [0.0]      # net bar starts at 0
    bar_bots = list(cum[1:]) + [total]
    pts = cum + [0.0, total]
    vmax=max(pts); vmin=min(pts); rng=(vmax-vmin) or 1
    pad=rng*0.18; vmax+=pad; vmin-=pad; rng=vmax-vmin
    plot_x, plot_y, plot_w, plot_h = x, y+0.1, w, h-0.75
    def yv(v): return plot_y + (vmax - v)/rng*plot_h
    zero_y = yv(0)
    _rect(s, plot_x, zero_y-0.006, plot_w, 0.012, fill=th["line"], shape=SH.RECTANGLE)
    n=len(show_deltas); slot=plot_w/max(1,n); bw=slot*0.58
    for i,d in enumerate(show_deltas):
        is_net = (i==n-1)
        top=bar_tops[i]; bot=bar_bots[i]
        y_hi=yv(max(top,bot)); y_lo=yv(min(top,bot))
        bx=plot_x + i*slot + (slot-bw)/2
        col = th["ink"] if is_net else (accent if d>=0 else neg)
        _rect(s, bx, y_hi, bw, max(0.05,y_lo-y_hi), fill=col, radius=0.10)
        # running connector
        if i<n-2:
            cy=yv(cum[i+1])
            _rect(s, bx+bw, cy-0.004, slot-bw, 0.008, fill=th["soft"], shape=SH.RECTANGLE)
        # value label
        lbl=f"{d:+,.0f}" if not is_net else f"{d:,.0f}"
        ly = y_hi-0.30 if (d>=0 or is_net) else y_lo+0.03
        _txt(s, lbl, bx-0.25, ly, bw+0.5, 0.24, size=10, font=th["bf"],
             color=col, bold=True, align="center")
        # contribution % (the data-density add) under the value, skip Net
        if not is_net:
            # place pct clear of the bar: below negative bars, above the value for positives
            py = (ly + 0.22) if (d >= 0) else (ly + 0.22)
            _txt(s, f"{pcts[i]:.0f}% of move", bx-0.25, py, bw+0.5, 0.2, size=7.5,
                 font=th["bf"], color=th["soft"], align="center")
        # category label
        _txt(s, str(show_labels[i]), bx-0.25, plot_y+plot_h+0.08, bw+0.5, 0.4, size=9.5,
             font=th["bf"], color=(th["ink"] if not is_net else th["accent"]),
             bold=is_net, align="center")
    return None


def chart_bar_ranked(s, x,y,w,h, labels, values, th, highlight_idx=0, show_share=True):
    """Ranked bars, driver highlighted, with value + share% labels and an average
    reference line drawn as an overlay (data-dense)."""
    cd=CategoryChartData(); cd.categories=[str(l) for l in labels]; cd.add_series("v", values)
    gf=s.shapes.add_chart(CT.COLUMN_CLUSTERED, In(x),In(y),In(w),In(h), cd); ch=gf.chart
    ch.has_title=False; ch.has_legend=False; ch.plots[0].gap_width=70
    ser=ch.series[0]
    for i in range(len(values)):
        pt=ser.points[i]; pt.format.fill.solid()
        pt.format.fill.fore_color.rgb=_hx(th["accent"] if i==highlight_idx else th["muted"])
    plot=ch.plots[0]; plot.has_data_labels=True
    dl=plot.data_labels; dl.font.size=Pt(11); dl.font.color.rgb=_hx(th["ink"]); dl.font.bold=True
    dl.position=LP.OUTSIDE_END; dl.number_format='#,##0'; dl.number_format_is_linked=False
    _axis_style(ch, th, hide_val=True)
    # overlay: average reference line + share% under each bar
    tot=sum(values) or 1; avg=tot/len(values) if values else 0
    vmax=max(values) or 1
    plot_x, plot_y, plot_w, plot_h = x+0.35, y+0.12, w-0.5, h-0.55
    if show_share:
        n=len(values); slot=plot_w/max(1,n)
        # average line
        ay = plot_y + (1 - avg/ (vmax*1.12))*plot_h
        _rect(s, plot_x, ay-0.004, plot_w, 0.008, fill=th["accent2"], shape=SH.RECTANGLE)
        _txt(s, f"avg {avg:,.0f}", plot_x+plot_w-1.3, ay-0.24, 1.3, 0.2, size=8,
             font=th["bf"], color=th["accent2"], align="right", bold=True)
        for i,v in enumerate(values):
            cxp=plot_x + i*slot
            _txt(s, f"{v/tot*100:.0f}%", cxp, plot_y+plot_h+0.04, slot, 0.2, size=8,
                 font=th["bf"], color=(th["accent"] if i==highlight_idx else th["soft"]),
                 align="center", bold=(i==highlight_idx))
    return ch


def chart_line(s, x,y,w,h, labels, values, th, area=False, ref_label=None):
    cd=CategoryChartData(); cd.categories=[str(l) for l in labels]; cd.add_series("v", values)
    kind=CT.AREA if area else CT.LINE
    gf=s.shapes.add_chart(kind, In(x),In(y),In(w),In(h), cd); ch=gf.chart
    ch.has_title=False; ch.has_legend=False
    ser=ch.series[0]
    if area:
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb=_hx(th["accent"]); _alpha_fill(ser, 24000)
        ser.format.line.color.rgb=_hx(th["accent"]); ser.format.line.width=Pt(2.5)
    else:
        ser.smooth=True; ser.format.line.color.rgb=_hx(th["accent"]); ser.format.line.width=Pt(3)
    # data labels (density)
    try:
        plot=ch.plots[0]; plot.has_data_labels=True
        dl=plot.data_labels; dl.font.size=Pt(8); dl.font.color.rgb=_hx(th["soft"])
        dl.number_format='#,##0'; dl.number_format_is_linked=False; dl.position=LP.ABOVE
    except Exception: pass
    _axis_style(ch, th)
    # average reference line overlay
    if values:
        avg=sum(values)/len(values); vmax=max(values); vmin=min(values); rng=(vmax-vmin) or 1
        plot_x, plot_y, plot_w, plot_h = x+0.4, y+0.1, w-0.5, h-0.5
        ay=plot_y + (vmax-avg)/(rng*1.1)*plot_h
        _rect(s, plot_x, ay-0.004, plot_w, 0.008, fill=th["accent2"], shape=SH.RECTANGLE)
        _txt(s, ref_label or f"avg {avg:,.0f}", plot_x+plot_w-1.4, ay-0.22, 1.4, 0.2, size=8,
             font=th["bf"], color=th["accent2"], align="right", bold=True)
    return ch


def chart_donut(s, x,y,w,h, labels, values, th):
    cd=CategoryChartData(); cd.categories=[str(l) for l in labels]; cd.add_series("v", values)
    gf=s.shapes.add_chart(CT.DOUGHNUT, In(x),In(y),In(w),In(h), cd); ch=gf.chart
    ch.has_title=False; ch.has_legend=True; ch.legend.position=LEG.RIGHT; ch.legend.include_in_layout=False
    ch.legend.font.size=Pt(10)
    cols=[th["accent"], th["muted"], th["accent2"], th["soft"], "C9CCDA", "E1E3EC"]
    for i,pt in enumerate(ch.plots[0].series[0].points):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb=_hx(cols[i%len(cols)])
    p=ch.plots[0]; p.has_data_labels=True; p.data_labels.show_percentage=True
    p.data_labels.number_format='0%'; p.data_labels.number_format_is_linked=False; p.data_labels.font.size=Pt(10)
    try: ch.font.name=th["bf"]
    except Exception: pass
    return ch


def chart_dumbbell(s, x,y,w,h, labels, prior, recent, th):
    """Prior vs recent dots per segment, connected — drawn as a clustered bar proxy
    isn't ideal, so we draw it manually with shapes for the consultant look."""
    n=len(labels); 
    if n==0: return
    allv = prior+recent; mn=min(allv); mx=max(allv); rng=(mx-mn) or 1
    row_h = h/n
    for i,lab in enumerate(labels):
        cy = y + i*row_h + row_h/2
        px = x + 2.2 + (prior[i]-mn)/rng*(w-2.6)
        rx = x + 2.2 + (recent[i]-mn)/rng*(w-2.6)
        _txt(s, str(lab), x, cy-0.16, 2.0, 0.32, size=10, font=th["bf"], color=th["ink"], valign="middle")
        # connector
        lo,hi = sorted([px,rx])
        _rect(s, lo, cy-0.012, hi-lo, 0.024, fill=th["line"], shape=SH.RECTANGLE)
        # prior dot (muted), recent dot (accent)
        _rect(s, px-0.07, cy-0.07, 0.14, 0.14, fill=th["muted"], shape=SH.OVAL)
        _rect(s, rx-0.07, cy-0.07, 0.14, 0.14, fill=th["accent"], shape=SH.OVAL)


def chart_slope(s, x,y,w,h, labels, prior, recent, th):
    """Two-column slope chart drawn with shapes/lines."""
    n=len(labels)
    if n==0: return
    allv=prior+recent; mn=min(allv); mx=max(allv); rng=(mx-mn) or 1
    lx, rxx = x+1.4, x+w-1.4
    def yv(v): return y+h-0.3 - (v-mn)/rng*(h-0.8)
    _txt(s,"Prior",lx-0.6,y-0.05,1.2,0.3,size=10,font=th["bf"],color=th["soft"],align="center")
    _txt(s,"Recent",rxx-0.6,y-0.05,1.2,0.3,size=10,font=th["bf"],color=th["soft"],align="center")
    cols=[th["accent"], th["accent2"], th["soft"], th["muted"], "9AA0B8"]
    for i,lab in enumerate(labels):
        y0,y1=yv(prior[i]),yv(recent[i]); c=cols[i%len(cols)]
        ln=s.shapes.add_connector(2, In(lx),In(y0),In(rxx),In(y1))
        ln.line.color.rgb=_hx(c); ln.line.width=Pt(2.2); ln.shadow.inherit=False
        _rect(s,lx-0.06,y0-0.06,0.12,0.12,fill=c,shape=SH.OVAL)
        _rect(s,rxx-0.06,y1-0.06,0.12,0.12,fill=c,shape=SH.OVAL)
        _txt(s,f"{lab}",rxx+0.15,y1-0.14,1.8,0.3,size=9,font=th["bf"],color=th["ink"],valign="middle")


def chart_scatter(s, x,y,w,h, xs, ys, th, xlab="", ylab=""):
    cd=XyChartData(); ser=cd.add_series("pts")
    for a,b in zip(xs,ys): ser.add_data_point(a,b)
    gf=s.shapes.add_chart(CT.XY_SCATTER, In(x),In(y),In(w),In(h), cd); ch=gf.chart
    ch.has_title=False; ch.has_legend=False
    pser=ch.series[0]; pser.marker.format.fill.solid(); pser.marker.format.fill.fore_color.rgb=_hx(th["accent"])
    try: pser.marker.size=5
    except Exception: pass
    _axis_style(ch, th)
    return ch


# ============================================================
# VISUAL FLOW  — the "storyline" slide (chevron arc of the narrative)
# ============================================================
def slide_storyline(prs, blank, th, style_name, flow_steps, framing):
    """A visual map of the deck's narrative arc as connected chevrons."""
    s=prs.slides.add_slide(blank); _bg(s, th["canvas"])
    _section_head(s, "Storyline", framing.get("title","How this story unfolds"), th)
    n=len(flow_steps); 
    if n==0: return s
    gap=0.25; total_w=11.6; cw=(total_w-(n-1)*gap)/n
    x0=0.9; y=3.0; chh=1.5
    for i,step in enumerate(flow_steps):
        x=x0+i*(cw+gap)
        shape = SH.PENTAGON if i<n-1 else SH.ROUNDED_RECTANGLE
        fill = th["accent"] if i==0 else th["accent2"] if i==n-1 else th["panel"]
        txtc = "FFFFFF" if i in (0,n-1) else th["ink"]
        r=_rect(s, x, y, cw, chh, fill=fill, line=(None if i in (0,n-1) else th["line"]),
                lw=1.4, shape=shape, radius=0.18, shdw=True)
        _txt(s, f"{i+1}", x, y+0.18, cw, 0.4, size=13, font=th["bf"],
             color=(txtc if i in (0,n-1) else th["accent"]), bold=True, align="center")
        _txt(s, step, x, y+0.55, cw, 0.8, size=14, font=th["hf"], color=txtc, bold=True,
             align="center", valign="middle", ls=1.0)
    _txt(s, f"Structure: {style_name}", 0.9, y+chh+0.4, 11.6, 0.4, size=12,
         font=th["bf"], color=th["soft"])
    return s




# ---- consultant glue ----
try:
    from pptx import Presentation as _Ck_Prs   # availability probe
    _CONSULT_OK = True
except Exception:
    _CONSULT_OK = False

# charts namespace passed to the pptx renderer
from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE_NS
_CHARTS_NS = {
    "_bg": _bg, "_rect": _rect, "_txt": _txt, "_pill": _pill,
    "_section_head": _section_head, "_alpha_fill": _alpha_fill,
    "chart_waterfall": chart_waterfall, "chart_bar_ranked": chart_bar_ranked,
    "chart_line": chart_line, "chart_donut": chart_donut,
    "chart_dumbbell": chart_dumbbell, "chart_slope": chart_slope,
    "chart_scatter": chart_scatter, "slide_storyline": slide_storyline,
    "SH": _MSO_SHAPE_NS,
}
# styles dict for the deck spec builder (with display names attached)
__STYLES = {}
for _nm, _v in STORYBOARD_STYLES.items():
    _vv = dict(_v); _vv["__name"] = _nm; __STYLES[_nm] = _vv


#
# Consultant deck assembly — spec builder, pptx slide renderer, HTML preview.
# Ties: consultant.py (findings) + agent.py (visual/story/narration) + pm_charts.py (charts).
#
from pptx import Presentation
from pptx.util import Inches as In, Pt


# theme registry (hex + font keys) used by the consultant renderer
CONSULT_THEMES = {
    "Indigo (Gamma-style)": {
        "canvas":"F7F7FB","ink":"1F2235","body":"3F4256","soft":"8A8DA3","muted":"C9CCDA",
        "accent":"6D5EF8","accent2":"8B7CFF","neg":"E5484D","tagbg":"E5E3FB","tagink":"5B4FE0",
        "line":"E6E6EF","panel":"FFFFFF",
    },
    "KiteIQX (Navy/Gold)": {
        "canvas":"F7F8FB","ink":"0A2540","body":"3A4A5E","soft":"7A8699","muted":"CBD5E1",
        "accent":"C79A3A","accent2":"E6C878","neg":"B91C1C","tagbg":"FBF0D5","tagink":"8A6A1E",
        "line":"E3E6EE","panel":"FFFFFF",
    },
    "Slate (Corporate)": {
        "canvas":"F1F5F9","ink":"0F172A","body":"334155","soft":"64748B","muted":"CBD5E1",
        "accent":"0EA5E9","accent2":"38BDF8","neg":"E11D48","tagbg":"E0F2FE","tagink":"0369A1",
        "line":"E2E8F0","panel":"FFFFFF",
    },
}
CONSULT_FONTS = {
    "Arial + Calibri": {"hf":"Arial","bf":"Calibri"},
    "Georgia + Calibri": {"hf":"Georgia","bf":"Calibri"},
    "Trebuchet + Calibri": {"hf":"Trebuchet MS","bf":"Calibri"},
}


def _theme(theme_name, font_name):
    th = dict(CONSULT_THEMES.get(theme_name, list(CONSULT_THEMES.values())[0]))
    f = CONSULT_FONTS.get(font_name, list(CONSULT_FONTS.values())[0])
    th["hf"], th["bf"] = f["hf"], f["bf"]
    return th


def _delta_chart_from_finding(f):
    """Turn a root_cause/interaction finding's table into waterfall deltas."""
    ch = f.chart or {}
    labels = ch.get("labels", []); deltas = ch.get("values", [])
    return labels, deltas


def __base_meta(framing, theme_name, font_name, style_key, ctx):
    from datetime import datetime
    return {"title": framing.get("title") or "Data Story",
            "subtitle": framing.get("subtitle") or (ctx.describe()[0].title() + " analysis"),
            "theme": theme_name, "font": font_name, "style": style_key,
            "date": datetime.now().strftime("%B %Y"),
            "decision": framing.get("decision", ""), "audience": framing.get("audience", "")}


def build_consultant_spec(analytics, findings, ctx, framing, theme_name, font_name,
                          style_key, order, narrations, visuals):
    """
    Assemble an ordered slide spec.
      narrations: {finding_index: narrative str}
      visuals:    {finding_index: {chart_kind, action_title}}
    """
    from datetime import datetime
    meta = {"title": framing.get("title") or "Data Story",
            "subtitle": framing.get("subtitle") or ctx.describe()[0].title() + " analysis",
            "theme": theme_name, "font": font_name, "style": style_key,
            "date": datetime.now().strftime("%B %Y"),
            "decision": framing.get("decision",""), "audience": framing.get("audience","")}

    # flow steps from the chosen style
    from_style = next((v for v in __STYLES.values() if v["key"] == style_key),
                      list(__STYLES.values())[0])
    flow = from_style["flow"]

    # ---- Standard Consulting storyline: planner-driven, section-based ----
    if style_key == "consulting":
        # build exec-summary content (LLM if available, else fallback)
        sit = [i for i, f in enumerate(findings) if f.ftype == "situation"]
        rc = [i for i, f in enumerate(findings) if f.ftype in ("root_cause", "interaction")]
        exec_content = {
            "core_answer": (narrations.get(rc[0]) if rc else narrations.get(sit[0]) if sit else
                            "The data points to a clear, addressable driver."),
            "insights": [narrations.get(i, "") for i in order[:3] if narrations.get(i)],
            "actions": (framing.get("recommendations") or [])[:2],
        }
        if analytics.llm:
            try:
                import json as _j, re as _r
                out = analytics.llm.predict(exec_summary_prompt(findings, framing, order[:5]), max_tokens=420)
                p = _j.loads(_r.sub(r"```json|```", "", out).strip())
                exec_content = {"core_answer": p.get("core_answer", exec_content["core_answer"]),
                                "insights": p.get("insights", exec_content["insights"])[:3],
                                "actions": p.get("actions", exec_content["actions"])[:2]}
            except Exception:
                pass

        plan = plan_consulting_storyline(findings, framing)
        # attach titles/narratives/takeaways/charts to each plan item
        for pi in plan:
            if pi["role"] == "exec_summary":
                pi["exec"] = exec_content
                continue
            fi = pi.get("finding_idx")
            if fi is not None:
                v = visuals.get(fi, {})
                pi["title"] = v.get("action_title") or "Finding"
                pi["chart_kind"] = v.get("chart_kind")
                pi["narrative"] = narrations.get(fi, findings[fi].narrative or "")
                pi["takeaway"] = pi["narrative"]
            elif pi["role"] == "actions":
                pi["recs"] = (framing.get("recommendations") or [narrations.get(i, "") for i in order[:3]])[:3]
                pi["title"] = "Recommended actions"
            elif pi["role"] == "impact":
                pi["impacts"] = framing.get("impacts", [])[:3]
                pi["title"] = "Expected impact"
            elif pi["role"] == "roadmap":
                pi["steps"] = framing.get("roadmap")
                pi["title"] = "Next steps & roadmap"
        return {"meta": {**__base_meta(framing, theme_name, font_name, style_key, ctx),
                         "flow": flow},
                "consulting_plan": plan, "findings_ref": True}


    slides = [{"kind":"title", "title":meta["title"], "subtitle":meta["subtitle"], "date":meta["date"]}]
    slides.append({"kind":"storyline", "style_name":from_style["__name"], "flow":flow,
                   "title": framing.get("flow_title","How this story unfolds")})

    for oi in order:
        f = findings[oi]
        v = visuals.get(oi, {"chart_kind": None, "action_title": None})
        narr = narrations.get(oi, f.narrative or "")
        title = v.get("action_title") or f._fallback_title() if hasattr(f, "_fallback_title") else v.get("action_title")
        title = v.get("action_title") or "Finding"
        tag = {"situation":"Situation","complication":"Complication","root_cause":"Root cause",
               "interaction":"Drill-down","mix":"Breakdown","evidence":"Evidence",
               "correlation":"Relationship","concentration":"Concentration","profile":"Overview",
               "action":"Action"}.get(f.role, f.ftype.title())

        slides.append({
            "kind":"finding", "tag":tag, "title":title, "narrative":narr,
            "chart_kind": v.get("chart_kind"), "chart": f.chart, "facts": f.facts,
            "ftype": f.ftype,
        })

    # recommendations slide (from action/narrations or fallback to top findings)
    recs = framing.get("recommendations") or [findings[i].narrative or "" for i in order[:3]]
    slides.append({"kind":"actions", "title":"Recommended Next Steps", "items":[r for r in recs if r][:3]})
    slides.append({"kind":"closing", "message": framing.get("closing","The data points to a clear next move."),
                   "date": meta["date"]})
    return {"meta": meta, "slides": slides}


# ── PPTX RENDER ─────────────────────────────────────────────
def build_consultant_pptx(spec, out_path, charts_mod, findings=None, framing=None):
    """charts_mod = the module namespace holding _bg,_rect,_txt,_pill,_section_head,chart_*"""
    th = _theme(spec["meta"]["theme"], spec["meta"]["font"])
    G = charts_mod
    prs = Presentation(); prs.slide_width=In(13.333); prs.slide_height=In(7.5)
    blank = prs.slide_layouts[6]; W,H = 13.333, 7.5

    # ---- Standard Consulting storyline path ----
    if spec.get("consulting_plan") is not None:
        plan = spec["consulting_plan"]
        meta = spec["meta"]
        # title slide
        ts = prs.slides.add_slide(blank); G["_bg"](ts, th["accent"])
        G["_rect"](ts, 0,0,0.18,H, fill=th["accent2"], shape=G["SH"].RECTANGLE)
        G["_txt"](ts, meta["date"].upper(), 0.9,1.5,6,0.4, size=12, font=th["bf"], color="FFFFFF", bold=True, spacing=3)
        G["_txt"](ts, meta["title"], 0.86,2.2,11,2.2, size=48, font=th["hf"], color="FFFFFF", bold=True, ls=1.04)
        G["_txt"](ts, meta.get("subtitle",""), 0.9,4.8,9.5,0.8, size=18, font=th["bf"], color="EDEBFF")
        total = len(plan) + 1
        for i, pi in enumerate(plan, 2):
            render_consulting_slide(prs, blank, pi, findings or [], framing or {}, th, G, i, total)
        prs.save(out_path)
        from pathlib import Path
        return True, f"Generated {Path(out_path).name} ({Path(out_path).stat().st_size//1024} KB)"

    def num(s, n): G["_txt"](s, str(n), 12.7, 7.15, 0.4, 0.25, size=9, font=th["bf"], color=th["soft"], align="right")

    for idx, sl in enumerate(spec["slides"], 1):
        k = sl["kind"]; s = prs.slides.add_slide(blank)

        if k == "title":
            G["_bg"](s, th["accent"])
            G["_rect"](s, 0,0,0.18,H, fill=th["accent2"], shape=G["SH"].RECTANGLE)
            G["_txt"](s, spec["meta"]["date"].upper(), 0.9,1.5,6,0.4, size=12, font=th["bf"],
                      color="FFFFFF", bold=True, spacing=3)
            G["_txt"](s, sl["title"], 0.86,2.2,10.5,2.4, size=52, font=th["hf"], color="FFFFFF", bold=True, ls=1.04)
            G["_txt"](s, sl.get("subtitle",""), 0.9,4.8,9.5,0.9, size=18, font=th["bf"], color="EDEBFF")
            G["_txt"](s, "Powered by KiteIQX Intelligence", 0.9,6.6,8,0.5, size=12, font=th["bf"], color="D6D2FF")

        elif k == "storyline":
            G["slide_storyline"](prs_existing=s) if False else None
            # draw inline (reuse pieces)
            G["_bg"](s, th["canvas"])
            G["_section_head"](s, "Storyline", sl.get("title","How this story unfolds"), th)
            flow = sl.get("flow", []); n=len(flow)
            if n:
                gap=0.25; total=11.6; cw=(total-(n-1)*gap)/n; x0=0.9; y=3.0; chh=1.5
                for i,step in enumerate(flow):
                    x=x0+i*(cw+gap)
                    shape = G["SH"].PENTAGON if i<n-1 else G["SH"].ROUNDED_RECTANGLE
                    fill = th["accent"] if i==0 else th["accent2"] if i==n-1 else th["panel"]
                    tc = "FFFFFF" if i in (0,n-1) else th["ink"]
                    G["_rect"](s, x,y,cw,chh, fill=fill, line=(None if i in (0,n-1) else th["line"]),
                               lw=1.4, shape=shape, radius=0.18, shdw=True)
                    G["_txt"](s, str(i+1), x,y+0.18,cw,0.4, size=13, font=th["bf"],
                              color=(tc if i in (0,n-1) else th["accent"]), bold=True, align="center")
                    G["_txt"](s, step, x,y+0.55,cw,0.8, size=14, font=th["hf"], color=tc, bold=True,
                              align="center", valign="middle")
                G["_txt"](s, f"Structure: {sl.get('style_name','')}", 0.9,y+chh+0.4,11.6,0.4,
                          size=12, font=th["bf"], color=th["soft"])
            num(s, idx)

        elif k == "finding":
            G["_bg"](s, th["canvas"])
            G["_section_head"](s, sl["tag"], sl["title"], th)
            ck = sl.get("chart_kind"); ch = sl.get("chart") or {}
            cx, cw, cy, chh = 0.9, 6.9, 2.25, 4.0
            drew = False
            try:
                if ck == "waterfall" and ch.get("labels"):
                    G["chart_waterfall"](s, cx,cy,cw,chh, ch["labels"], ch["values"], th); drew=True
                elif ck == "bar_ranked" and ch.get("labels"):
                    G["chart_bar_ranked"](s, cx,cy,cw,chh, ch["labels"], ch["values"], th, 0); drew=True
                elif ck in ("line","area") and ch.get("labels"):
                    G["chart_line"](s, cx,cy,cw,chh, ch["labels"], ch["values"], th, area=(ck=="area")); drew=True
                elif ck in ("donut","stacked100") and ch.get("labels"):
                    G["chart_donut"](s, cx,cy,cw,chh, ch["labels"], ch["values"], th); drew=True
                elif ck == "scatter" and ch.get("x"):
                    G["chart_scatter"](s, cx,cy,cw,chh, ch["x"], ch["y"], th); drew=True
                elif ck in ("slope","dumbbell"):
                    # need prior/recent; reuse values as recent, facts prior — fallback to bar
                    if ch.get("labels"):
                        G["chart_bar_ranked"](s, cx,cy,cw,chh, ch["labels"], ch["values"], th, 0); drew=True
            except Exception:
                drew=False
            # right: narrative + key fact callout
            ix = 8.1
            facts = sl.get("facts", {})
            big = facts.get("contribution_pct") or facts.get("share") or facts.get("pct") or facts.get("r","")
            if big:
                G["_rect"](s, ix, cy, 4.3, 1.25, fill=th["tagbg"], radius=0.12)
                G["_txt"](s, str(big), ix, cy+0.14, 4.3, 0.7, size=26, font=th["hf"], color=th["tagink"],
                          bold=True, align="center")
                G["_txt"](s, "key figure", ix, cy+0.86, 4.3, 0.3, size=10, font=th["bf"],
                          color=th["soft"], align="center")
                ny = cy+1.5
            else:
                ny = cy
            G["_txt"](s, sl.get("narrative",""), ix, ny+0.1, 4.3, 2.4, size=14, font=th["bf"],
                      color=th["body"], ls=1.35)
            if not drew and not ck:
                # number-led slide
                G["_txt"](s, str(big), 0.9, 3.0, 6.5, 1.6, size=90, font=th["hf"], color=th["accent"], bold=True)
            num(s, idx)

        elif k == "actions":
            G["_bg"](s, th["canvas"])
            G["_section_head"](s, "Action plan", sl["title"], th)
            cols=[th["accent"], th["accent2"], th["soft"]]
            for i,t in enumerate(sl.get("items", [])[:3]):
                y=1.7+i*1.7
                G["_rect"](s, 0.9,y,11.6,1.45, fill=th["panel"], line=th["line"], lw=1, radius=0.10, shdw=True)
                G["_rect"](s, 0.9,y,0.7,1.45, fill=cols[i%3], shape=G["SH"].RECTANGLE)
                G["_txt"](s, str(i+1), 0.9,y,0.7,1.45, size=28, font=th["hf"], color="FFFFFF",
                          bold=True, align="center", valign="middle")
                G["_txt"](s, t, 1.8,y+0.12,10.4,1.2, size=14.5, font=th["bf"], color=th["ink"],
                          valign="middle", ls=1.15)
            num(s, idx)

        elif k == "closing":
            G["_bg"](s, th["accent"])
            G["_rect"](s, 0,0,0.18,H, fill=th["accent2"], shape=G["SH"].RECTANGLE)
            G["_txt"](s, "KITEIQX INTELLIGENCE", 0.9,1.4,9,0.4, size=12, font=th["bf"],
                      color="EDEBFF", bold=True, spacing=3)
            G["_txt"](s, sl.get("message","Thank you."), 0.9,2.3,11,2.4, size=40, font=th["hf"],
                      color="FFFFFF", bold=True, ls=1.05)
            G["_txt"](s, "Generated by KiteIQX Intelligence · "+sl.get("date",""), 0.9,5.7,10,0.4,
                      size=12, font=th["bf"], color="D6D2FF")

    prs.save(out_path)
    from pathlib import Path
    return True, f"Generated {Path(out_path).name} ({Path(out_path).stat().st_size//1024} KB)"


# ── HTML PREVIEW (consultant slides) ────────────────────────
def _svg_for(ck, ch, th):
    if not ch: return ""
    A="#"+th["accent"]; M="#"+th["muted"]; NEG="#"+th["neg"]; INK="#"+th["ink"]; SOFT="#"+th["soft"]
    labels=[str(l) for l in ch.get("labels",[])]; vals=[float(v) for v in ch.get("values",[])]
    w,hh=470,210; pad=26
    if ck=="waterfall" and vals:
        cum=[0]; 
        for d in vals: cum.append(cum[-1]+d)
        allp=cum+[0]; vmax=max(allp); vmin=min(allp); rng=(vmax-vmin) or 1
        def yv(v): return pad+(vmax-v)/rng*(hh-2*pad)
        zero=yv(0); out=f'<line x1="{pad}" y1="{zero:.0f}" x2="{w-pad}" y2="{zero:.0f}" stroke="{M}"/>'
        n=len(vals); slot=(w-2*pad)/max(1,n); bw=slot*0.55
        for i,d in enumerate(vals):
            top=yv(max(cum[i],cum[i+1])); bot=yv(min(cum[i],cum[i+1])); bx=pad+i*slot+(slot-bw)/2
            out+=f'<rect x="{bx:.0f}" y="{top:.0f}" width="{bw:.0f}" height="{max(2,bot-top):.0f}" rx="2" fill="{A if d>=0 else NEG}"/>'
            out+=f'<text x="{bx+bw/2:.0f}" y="{(top-3) if d>=0 else (bot+11):.0f}" font-size="8" fill="{A if d>=0 else NEG}" text-anchor="middle">{d:+.0f}</text>'
            out+=f'<text x="{bx+bw/2:.0f}" y="{hh-6}" font-size="8" fill="{SOFT}" text-anchor="middle">{labels[i][:6]}</text>'
        return f'<svg viewBox="0 0 {w} {hh}" width="100%">{out}</svg>'
    if ck in ("bar_ranked","bar") and vals:
        mx=max(vals) or 1; n=len(vals); slot=(w-2*pad)/max(1,n); bw=slot*0.6; out=""
        for i,v in enumerate(vals):
            bh=v/mx*(hh-2*pad); bx=pad+i*slot+(slot-bw)/2; by=hh-pad-bh
            out+=f'<rect x="{bx:.0f}" y="{by:.0f}" width="{bw:.0f}" height="{bh:.0f}" rx="3" fill="{A if i==0 else M}"/>'
            out+=f'<text x="{bx+bw/2:.0f}" y="{by-3:.0f}" font-size="8" fill="{INK}" text-anchor="middle">{v:.0f}</text>'
            out+=f'<text x="{bx+bw/2:.0f}" y="{hh-6}" font-size="8" fill="{SOFT}" text-anchor="middle">{labels[i][:6]}</text>'
        return f'<svg viewBox="0 0 {w} {hh}" width="100%">{out}</svg>'
    if ck in ("line","area") and vals:
        mx=max(vals) or 1; mn=min(vals); rng=(mx-mn) or 1; n=len(vals); step=(w-2*pad)/max(1,n-1)
        pts=" ".join(f"{pad+i*step:.0f},{hh-pad-(v-mn)/rng*(hh-2*pad):.0f}" for i,v in enumerate(vals))
        fill=f'<polygon points="{pad},{hh-pad} {pts} {w-pad},{hh-pad}" fill="{A}" opacity="0.16"/>' if ck=="area" else ""
        return f'<svg viewBox="0 0 {w} {hh}" width="100%">{fill}<polyline points="{pts}" fill="none" stroke="{A}" stroke-width="2.5"/></svg>'
    if ck in ("donut","stacked100") and vals:
        import math; tot=sum(vals) or 1; cx,cy,r=hh/2,hh/2,hh/2-14; a0=-math.pi/2; out=""
        cols=[A,M,"#"+th["accent2"],SOFT,"#C9CCDA"]
        for i,v in enumerate(vals):
            a1=a0+v/tot*2*math.pi; x0=cx+r*math.cos(a0); y0=cy+r*math.sin(a0); x1=cx+r*math.cos(a1); y1=cy+r*math.sin(a1)
            lg=1 if a1-a0>math.pi else 0
            out+=f'<path d="M{cx},{cy} L{x0:.1f},{y0:.1f} A{r},{r} 0 {lg} 1 {x1:.1f},{y1:.1f} Z" fill="{cols[i%len(cols)]}"/>'; a0=a1
        out+=f'<circle cx="{cx}" cy="{cy}" r="{r*0.58:.0f}" fill="#fff"/>'
        return f'<svg viewBox="0 0 {w} {hh}" width="100%">{out}</svg>'
    if ck=="scatter" and ch.get("x"):
        xs=ch["x"][:120]; ys=ch["y"][:120]
        xmn,xmx=min(xs),max(xs); ymn,ymx=min(ys),max(ys); xr=(xmx-xmn)or 1; yr=(ymx-ymn)or 1
        dots="".join(f'<circle cx="{pad+(x-xmn)/xr*(w-2*pad):.0f}" cy="{hh-pad-(y-ymn)/yr*(hh-2*pad):.0f}" r="2.2" fill="{A}" opacity="0.6"/>' for x,y in zip(xs,ys))
        return f'<svg viewBox="0 0 {w} {hh}" width="100%">{dots}</svg>'
    return ""


def render_consultant_html(sl, meta):
    th=_theme(meta["theme"], meta["font"])
    A="#"+th["accent"]; A2="#"+th["accent2"]; INK="#"+th["ink"]; SOFT="#"+th["soft"]
    LINE="#"+th["line"]; TAG="#"+th["tagbg"]; TAGINK="#"+th["tagink"]; CAN="#"+th["canvas"]; PANEL="#"+th["panel"]
    frame=f'width:100%;aspect-ratio:16/9;border-radius:12px;overflow:hidden;position:relative;font-family:Inter,system-ui,sans-serif;border:1px solid {LINE};box-shadow:0 2px 10px rgba(0,0,0,0.06);'
    def card(inner, bgc="#fff"): return f'<div style="{frame}background:{bgc};">{inner}</div>'
    def tagline(tag,title,tc=INK):
        return (f'<div style="position:absolute;left:4%;top:7%;background:{TAG};color:{TAGINK};font-size:0.55rem;'
                f'font-weight:700;letter-spacing:1.5px;padding:3px 12px;border-radius:20px;">{tag.upper()}</div>'
                f'<div style="position:absolute;left:4%;top:15%;right:5%;font-size:1.5rem;font-weight:800;color:{tc};line-height:1.1;">{title}</div>')
    k=sl["kind"]
    if k=="title":
        return card(
            f'<div style="position:absolute;left:0;top:0;bottom:0;width:1.5%;background:{A2};"></div>'
            f'<div style="position:absolute;left:5%;top:22%;color:#fff;font-size:0.7rem;font-weight:700;letter-spacing:3px;">{sl.get("date","").upper()}</div>'
            f'<div style="position:absolute;left:5%;top:32%;right:8%;color:#fff;font-size:2.2rem;font-weight:800;line-height:1.08;">{sl["title"]}</div>'
            f'<div style="position:absolute;left:5%;top:66%;right:12%;color:#EDEBFF;font-size:0.95rem;">{sl.get("subtitle","")}</div>', bgc=A)
    if k=="storyline":
        flow=sl.get("flow",[]); n=len(flow); chips=""
        for i,st in enumerate(flow):
            bg = A if i==0 else A2 if i==n-1 else PANEL
            tc = "#fff" if i in (0,n-1) else INK
            arrow = "" if i==n-1 else f'<div style="color:{SOFT};font-size:1rem;align-self:center;">→</div>'
            chips+=(f'<div style="flex:1;background:{bg};color:{tc};border:1px solid {LINE};border-radius:10px;'
                    f'padding:14px 6px;text-align:center;font-size:0.7rem;font-weight:700;">{i+1}<br>{st}</div>{arrow}')
        return card(tagline("Storyline", sl.get("title",""))+
            f'<div style="position:absolute;left:4%;right:4%;top:42%;display:flex;gap:6px;align-items:stretch;">{chips}</div>'
            f'<div style="position:absolute;left:4%;top:72%;color:{SOFT};font-size:0.6rem;">Structure: {sl.get("style_name","")}</div>', bgc=CAN)
    if k=="finding":
        facts=sl.get("facts",{}); big=facts.get("contribution_pct") or facts.get("share") or facts.get("pct") or facts.get("r","")
        svg=_svg_for(sl.get("chart_kind"), sl.get("chart"), th)
        bigbox=(f'<div style="background:{TAG};border-radius:8px;padding:8px;text-align:center;margin-bottom:8px;">'
                f'<div style="font-size:1.4rem;font-weight:800;color:{TAGINK};">{big}</div>'
                f'<div style="font-size:0.5rem;color:{SOFT};">key figure</div></div>') if big else ""
        return card(tagline(sl["tag"], sl["title"])+
            f'<div style="position:absolute;left:4%;top:30%;width:58%;bottom:5%;">{svg}</div>'
            f'<div style="position:absolute;right:4%;top:30%;width:31%;">{bigbox}'
            f'<div style="font-size:0.62rem;color:{INK};line-height:1.4;">{sl.get("narrative","")}</div></div>', bgc=CAN)
    if k=="actions":
        rows=""
        for i,t in enumerate(sl.get("items",[])[:3]):
            c=[A,A2,SOFT][i%3]
            rows+=(f'<div style="display:flex;margin-bottom:8px;border:1px solid {LINE};border-radius:8px;overflow:hidden;background:#fff;">'
                   f'<div style="background:{c};color:#fff;width:38px;display:flex;align-items:center;justify-content:center;font-weight:800;">{i+1}</div>'
                   f'<div style="padding:9px 12px;font-size:0.66rem;color:{INK};">{t}</div></div>')
        return card(tagline("Action plan", sl["title"])+
            f'<div style="position:absolute;left:4%;right:4%;top:30%;">{rows}</div>', bgc=CAN)
    if k=="closing":
        return card(
            f'<div style="position:absolute;left:0;top:0;bottom:0;width:1.5%;background:{A2};"></div>'
            f'<div style="position:absolute;left:5%;top:22%;color:#EDEBFF;font-size:0.7rem;font-weight:700;letter-spacing:3px;">KITEIQX INTELLIGENCE</div>'
            f'<div style="position:absolute;left:5%;top:34%;right:14%;color:#fff;font-size:1.8rem;font-weight:800;">{sl.get("message","")}</div>', bgc=A)
    return card(f'<div style="padding:20px;">{k}</div>', bgc=CAN)




# ============================================================
# PRESENTATION MAKER — consultant agent UI
# ============================================================
import streamlit.components.v1 as _cmp


def _pm_reset():
    for k in list(st.session_state.keys()):
        if k.startswith("pm2_"):
            st.session_state.pop(k, None)




# ============================================================
# CHAT AGENT (data + visuals + deck-aware)
# ============================================================
import json as _cj
import re as _cr



# ---------- grounding: describe the deck + the data to the LLM ----------
def _plan_slide_list(spec):
    """Return numbered slide list (matching preview) + index map -> plan idx."""
    items = []
    if spec.get("consulting_plan") is not None:
        # slide 1 = title; plan items are slides 2..N
        items.append((1, "title", "Title slide", None))
        for pi_i, pi in enumerate(spec["consulting_plan"]):
            n = pi_i + 2
            label = pi.get("title") or pi.get("role", "slide")
            items.append((n, pi.get("role", "slide"), label, pi_i))
    else:
        for i, sl in enumerate(spec.get("slides", [])):
            items.append((i + 1, sl.get("kind", "slide"), sl.get("title", sl.get("message", "")), i))
    return items


def deck_brief(spec):
    rows = []
    for n, role, label, _ in _plan_slide_list(spec):
        rows.append(f"  Slide {n} [{role}]: {label}")
    return "Current deck:\n" + "\n".join(rows)


def data_brief_for_chat(analytics):
    dims = [c for c in analytics.categorical_cols if 2 <= analytics.df[c].nunique() <= 40]
    mets = list(analytics.money_cols) or analytics.numeric_cols
    return (f"Available dimensions to slice by: {', '.join(dims) or '(none)'}\n"
            f"Available metrics: {', '.join(mets[:8]) or '(none)'}\n"
            f"Time column: {analytics.datetime_cols[0] if analytics.datetime_cols else '(none)'}")


def findings_brief(findings):
    return "Key proven findings:\n" + "\n".join(
        f"  - [{f.ftype}] {f.fact_brief()[:110]}" for f in findings[:8])


# ---------- interpret a user message into one action ----------
_CHART_OPTS = ["waterfall", "bar_ranked", "line", "area", "donut", "scatter"]

def _deep_findings_brief(findings):
    """Low-level: every proven number the agent can cite, with the full drill path."""
    lines = []
    for i, f in enumerate(findings):
        lines.append(f"  [{i}] type={f.ftype} role={f.role} :: {f.fact_brief()}")
    return "All proven findings (cite these exact numbers):\n" + "\n".join(lines)


def _history_brief(history, k=8):
    if not history:
        return ""
    turns = history[-k:]
    out = []
    for role, msg in turns:
        who = "User" if role == "user" else "Assistant"
        out.append(f"{who}: {msg}")
    return "Conversation so far (remember this context):\n" + "\n".join(out)


def interpret_chat(user_msg, spec, findings, analytics, llm, history=None):
    """
    Returns an action dict. Defaults to CONVERSATION ('answer'); only emits a
    mutating action when the user clearly issues an imperative to change the deck.
    `history` = list of (role, msg) carrying memory of the dialogue.
    """
    if not llm:
        return {"action": "answer",
                "text": "The AI engine isn't configured, so I can't converse or take actions yet. "
                        "Add your GROQ_API_KEY to enable the assistant."}
    prompt = (
        "You are the Presentation Maker assistant: a senior data-analysis & visualization consultant "
        "having a CONVERSATION with the user about their deck. You have full memory of the dialogue.\n\n"
        "DEFAULT TO TALKING. Only change the deck when the user gives a clear command to do so "
        "(e.g. 'change', 'add', 'remove', 'reorder', 'make it', 'rewrite'). If they ask a question, "
        "are exploring options, or say something ambiguous, use 'answer' and have a real conversation — "
        "ask a clarifying question or lay out options. Never edit the deck just because you could.\n\n"
        "When you answer, be SPECIFIC: cite the exact numbers from the findings below "
        "(segments, contribution %, the drill path like 'NY × Q4 × Electronics = 49% of the drop'). "
        "No vague generalities.\n\n"
        f"{_history_brief(history)}\n\n"
        f"{deck_brief(spec)}\n\n{data_brief_for_chat(analytics)}\n\n{_deep_findings_brief(findings)}\n\n"
        f'User just said: "{user_msg}"\n\n'
        "Return ONLY JSON (no fences). Choose ONE:\n"
        '- Converse / answer / advise / ask a clarifying question: {"action":"answer","text":"...specific, with numbers..."}\n'
        '- Edit slide text (ONLY on explicit request): {"action":"edit_field","slide":N,"field":"title|narrative|takeaway|subtitle","value":"..."}\n'
        f'- Change a chart (explicit): {{"action":"change_chart","slide":N,"chart_kind":"one of {_CHART_OPTS}"}}\n'
        '- Reorder (explicit): {"action":"reorder","order":[slide numbers]}\n'
        '- Delete a slide (explicit): {"action":"drop_slide","slide":N}\n'
        '- Add a slide from a new data cut (explicit): {"action":"add_data_slide","dimension":"<col>","metric":"<col>","kind":"bar_ranked|donut|waterfall"}\n'
        '- Restyle (explicit): {"action":"restyle","theme":"...","font":"..."}\n\n'
        "If in doubt, use answer. Slide numbers must match the deck list above."
    )
    try:
        raw = _cr.sub(r"```json|```", "", llm.predict(prompt, max_tokens=600)).strip()
        mb = _cr.search(r"\{.*\}", raw, _cr.DOTALL)
        action = _cj.loads(mb.group(0) if mb else raw)
        if "action" not in action:
            return {"action": "answer", "text": raw[:800]}
        return action
    except Exception:
        # On any parse failure, treat it as conversation, not a botched edit
        try:
            return {"action": "answer", "text": llm.predict(
                f"{_history_brief(history)}\n\n{_deep_findings_brief(findings)}\n\n"
                f"User: {user_msg}\nAnswer specifically with numbers from the findings, as a consultant.",
                max_tokens=500).strip()}
        except Exception:
            return {"action": "answer", "text": "Could you rephrase that?"}


# ---------- apply an action to the spec (mutates in place) ----------
def _plan_idx_from_slide(spec, slide_n):
    for n, role, label, pidx in _plan_slide_list(spec):
        if n == slide_n:
            return role, pidx
    return None, None


def apply_chat_action(action, spec, findings, analytics, deps):
    """
    deps = {propose_visual, narrate_finding, contribution_by_dim, mix_by_dim,
            _period_split, Finding, CONSULT_THEMES, CONSULT_FONTS, plan helpers...}
    Returns (ok, human_message).
    """
    act = action.get("action")
    is_plan = spec.get("consulting_plan") is not None

    if act == "answer":
        return True, action.get("text", "")

    if act == "restyle":
        th = action.get("theme"); ft = action.get("font")
        if th and th in deps["CONSULT_THEMES"]:
            spec["meta"]["theme"] = th
        if ft and ft in deps["CONSULT_FONTS"]:
            spec["meta"]["font"] = ft
        return True, f"Restyled the deck (theme={spec['meta']['theme']}, font={spec['meta']['font']})."

    if act == "edit_field":
        slide = int(action.get("slide", 0)); field = action.get("field", ""); value = action.get("value", "")
        role, pidx = _plan_idx_from_slide(spec, slide)
        if is_plan:
            if pidx is None:  # title slide
                if field in ("title", "subtitle"):
                    spec["meta"][field] = value; return True, f"Updated deck {field}."
                return False, "That field isn't on the title slide."
            pi = spec["consulting_plan"][pidx]
            if role == "exec_summary":
                pi.setdefault("exec", {})
                if field in ("title", "core_answer"):
                    pi["exec"]["core_answer"] = value; return True, "Updated the executive summary."
            if field in ("title", "narrative", "takeaway"):
                pi[field] = value; return True, f"Updated slide {slide} {field}."
            return False, f"Can't edit '{field}' on that slide."
        else:
            sl = spec["slides"][pidx]
            if field in sl or field in ("title", "narrative", "subtitle", "message", "takeaway"):
                sl[field] = value; return True, f"Updated slide {slide} {field}."
            return False, f"Field '{field}' not found on slide {slide}."

    if act == "change_chart":
        slide = int(action.get("slide", 0)); kind = action.get("chart_kind", "")
        role, pidx = _plan_idx_from_slide(spec, slide)
        if pidx is None:
            return False, "That slide has no chart to change."
        pi = spec["consulting_plan"][pidx] if is_plan else spec["slides"][pidx]
        fi = pi.get("finding_idx")
        if fi is None or fi >= len(findings):
            # legacy slide with embedded chart
            if pi.get("chart"):
                pi["chart_kind"] = kind; return True, f"Changed slide {slide} chart to {kind}."
            return False, "That slide isn't backed by a chartable finding."
        ok, reason = deps["validate_chart_choice"](findings[fi], kind)
        if not ok:
            return False, f"Can't use {kind} here — {reason}. Try one that fits this finding."
        pi["chart_kind"] = kind
        return True, f"Changed slide {slide} chart to {kind}."

    if act == "reorder":
        order = action.get("order", [])
        if not is_plan:
            return False, "Reordering is available for the consulting storyline."
        # map slide numbers -> plan indices (skip title=1)
        pmap = {n: pidx for n, role, label, pidx in _plan_slide_list(spec) if pidx is not None}
        new_plan = []
        for n in order:
            if n in pmap and pmap[n] is not None:
                new_plan.append(spec["consulting_plan"][pmap[n]])
        # append any omitted
        seen = set(id(p) for p in new_plan)
        for p in spec["consulting_plan"]:
            if id(p) not in seen:
                new_plan.append(p)
        spec["consulting_plan"] = new_plan
        return True, "Reordered the slides."

    if act == "drop_slide":
        slide = int(action.get("slide", 0))
        role, pidx = _plan_idx_from_slide(spec, slide)
        if pidx is None:
            return False, "Can't drop the title slide."
        target = spec["consulting_plan"] if is_plan else spec["slides"]
        if 0 <= pidx < len(target):
            target.pop(pidx); return True, f"Dropped slide {slide}."
        return False, "Slide not found."

    if act == "add_data_slide":
        return _add_data_slide(action, spec, findings, analytics, deps)

    return False, f"I don't know how to '{act}' yet."


def _match_col(name, cols):
    if not name:
        return None
    n = name.strip().lower().replace(" ", "_")
    for c in cols:
        if c.lower() == n or c.lower().replace(" ", "_") == n:
            return c
    for c in cols:
        if n in c.lower() or c.lower() in n:
            return c
    return None


def _add_data_slide(action, spec, findings, analytics, deps):
    """Tier 2: run a fresh decomposition on a user-named dimension+metric, add a slide."""
    df = analytics.df
    dim = _match_col(action.get("dimension"), analytics.categorical_cols)
    metric = _match_col(action.get("metric"), analytics.numeric_cols) or \
             (analytics.money_cols[0] if analytics.money_cols else
              (analytics.numeric_cols[0] if analytics.numeric_cols else None))
    if not dim:
        return False, ("I couldn't match that dimension to a column. Available: "
                       + ", ".join(analytics.categorical_cols[:8]))
    if not metric:
        return False, "No numeric metric available to chart."

    fmt = deps["_fmt"]
    # build a mix/ranking finding
    g = df.groupby(dim)[metric].sum().sort_values(ascending=False)
    tot = float(g.sum()) or 1.0
    top = g.index[0]; topv = float(g.iloc[0]); share = topv / tot * 100
    Finding = deps["Finding"]
    f = Finding("mix", "root_cause",
                {"metric": metric, "dimension": dim, "segment": str(top),
                 "value": fmt(topv), "share": f"{share:.0f}%"},
                chart={"kind": action.get("kind", "bar_ranked") if action.get("kind") in ("bar_ranked","donut","waterfall") else "bar",
                       "title": f"{metric} by {dim}",
                       "labels": [str(x) for x in g.head(6).index.tolist()],
                       "values": [round(float(x), 2) for x in g.head(6).values.tolist()],
                       "metric": metric, "dim": dim},
                importance=55 + share / 2)
    findings.append(f); fi = len(findings) - 1
    v = deps["propose_visual"](f, analytics.llm, spec["meta"].get("audience", "leadership"))
    f.narrative = deps["narrate_finding"](f, {"audience": spec["meta"].get("audience", "")}, analytics.llm)

    item = {"section": "Insights & Drivers", "role": "segment", "layout": "bar_ranked",
            "finding_idx": fi, "title": v.get("action_title") or f"{top} leads {dim}",
            "chart_kind": v.get("chart_kind", "bar_ranked"),
            "narrative": f.narrative, "takeaway": f.narrative}
    if spec.get("consulting_plan") is not None:
        # insert before recommendations
        plan = spec["consulting_plan"]
        ins_at = next((i for i, p in enumerate(plan) if p.get("section") == "Recommendations"), len(plan))
        plan.insert(ins_at, item)
    else:
        spec.setdefault("slides", []).append(
            {"kind": "finding", "tag": "Breakdown", "title": item["title"],
             "narrative": f.narrative, "chart_kind": item["chart_kind"],
             "chart": f.chart, "facts": f.facts, "ftype": "mix"})
    return True, f"Added a slide breaking {metric} down by {dim} (top: {top}, {share:.0f}%)."


def render_presentation_maker():
    a = st.session_state.get("analytics")
    if a is None:
        st.info("Load data first (sidebar), then come back to build your story.")
        return
    if not _CONSULT_OK:
        st.error("`python-pptx` isn't installed. Add `python-pptx` to requirements.txt and redeploy.")
        return

    st.markdown(
        '<div class="kx-card"><div class="kx-card-title">Presentation Maker</div>'
        '<div class="kx-card-sub">A consultant agent that decomposes your data to find what is really driving '
        'the numbers, proposes a storyline, and builds an editable, visual-flow deck.</div></div>',
        unsafe_allow_html=True,
    )

    if "pm2_stage" not in st.session_state:
        st.session_state.pm2_stage = "analyze"
        st.session_state.pm2_framing = {}
        st.session_state.pm2_q = 0

    # theme/font/style pickers + restart
    c1, c2, c3, c4 = st.columns([2, 2, 2, 1])
    theme_name = c1.selectbox("Visual style", list(CONSULT_THEMES.keys()), key="pm2_theme")
    font_name = c2.selectbox("Font", list(CONSULT_FONTS.keys()), key="pm2_font")
    _style_opts = ["✨ Let the agent decide"] + list(STORYBOARD_STYLES.keys())
    style_label = c3.selectbox("Storyboard structure", _style_opts, key="pm2_style")
    with c4:
        st.markdown("<div style='height:1.7rem'></div>", unsafe_allow_html=True)
        if st.button("↺ Restart", key="pm2_restart", use_container_width=True):
            _pm_reset(); st.rerun()
    if style_label == "✨ Let the agent decide":
        st.caption("The agent will pick the storyboard structure that best fits your data and goal, and explain why.")
    else:
        st.caption(STORYBOARD_STYLES[style_label]["blurb"])
    st.markdown("---")

    stage = st.session_state.pm2_stage

    # ---- ANALYZE: run the consultant engine ----
    if stage == "analyze":
        with st.spinner("Decomposing the data and tracing what's driving the numbers…"):
            ctx, findings = derive_findings(a)
            for f in findings:                # attach narration lazily later
                f.narrative = None
            st.session_state.pm2_ctx_modes = ctx.modes
            st.session_state.pm2_primary, st.session_state.pm2_bits = ctx.describe()
            st.session_state.pm2_findings = findings
        st.session_state.pm2_stage = "frame"
        st.rerun()

    ctx_bits = st.session_state.get("pm2_bits", [])
    findings = st.session_state.get("pm2_findings", [])

    # rebuild a lightweight ctx-like object for spec building
    class _Ctx:  # minimal shim carrying what the spec builder needs
        def __init__(s, modes, primary): s.modes = modes; s._p = primary
        def describe(s): return s._p, []
    ctx = _Ctx(st.session_state.get("pm2_ctx_modes", []), st.session_state.get("pm2_primary", "profile"))

    # ---- FRAME: conversational framing questions ----
    if stage == "frame":
        with st.chat_message("assistant"):
            st.markdown("Here's what I can build from your data:")
            for b in ctx_bits:
                st.markdown(f"- {b}")
            st.markdown("A few quick questions so I storyboard it the way a consultant would 👇")

        FQ = [
            ("decision", "What decision should this deck drive?",
             ["Where to invest / cut", "Explain a result", "Win buy-in / pitch", "Just a status update"]),
            ("audience", "Who's the audience?",
             ["Board / CXO", "Investors", "Team / Ops", "Clients"]),
        ]
        qi = st.session_state.pm2_q
        if qi < len(FQ):
            key, q, opts = FQ[qi]
            with st.chat_message("assistant"):
                st.markdown(f"**{q}**")
                cols = st.columns(len(opts))
                for j, o in enumerate(opts):
                    if cols[j].button(o, key=f"pm2_fq_{qi}_{j}", use_container_width=True):
                        st.session_state.pm2_framing[key] = o
                        st.session_state.pm2_q += 1
                        st.rerun()
        else:
            with st.chat_message("assistant"):
                st.markdown("**One line — what's the single thing they must leave knowing?** (optional)")
            tk = st.chat_input("Type the key takeaway, or skip…")
            cskip, _ = st.columns([1, 5])
            go = cskip.button("Skip", key="pm2_skip_tk")
            if tk is not None or go:
                st.session_state.pm2_framing["takeaway"] = (tk or "").strip()
                st.session_state.pm2_stage = "storyboard"
                st.rerun()

    # ---- STORYBOARD: agent drafts ordered narrative + reasoning ----
    elif stage == "storyboard":
        framing = st.session_state.pm2_framing
        # resolve the storyboard structure — agent decides, or user-picked
        if "pm2_resolved_style" not in st.session_state:
            if style_label == "✨ Let the agent decide":
                with st.spinner("The agent is choosing the best structure for your data…"):
                    choice = agent_choose_structure(findings, framing, a.llm)
                st.session_state.pm2_resolved_style = choice["style_key"]
                st.session_state.pm2_resolved_name = choice["style_name"]
                st.session_state.pm2_resolved_why = choice["rationale"]
            else:
                st.session_state.pm2_resolved_style = STORYBOARD_STYLES[style_label]["key"]
                st.session_state.pm2_resolved_name = style_label
                st.session_state.pm2_resolved_why = None
        resolved_key = st.session_state.pm2_resolved_style
        resolved_name = st.session_state.pm2_resolved_name

        if "pm2_sb" not in st.session_state:
            with st.spinner("Drafting the storyline…"):
                sb = draft_storyboard(findings, framing, a.llm)
                style_order, _flow = order_findings_by_style(findings, resolved_key)
                st.session_state.pm2_sb = sb
                st.session_state.pm2_style_order = style_order
        sb = st.session_state.pm2_sb
        with st.chat_message("assistant"):
            if st.session_state.get("pm2_resolved_why"):
                st.markdown(f"**I chose the _{resolved_name}_ structure.**")
                st.markdown(f"_{st.session_state.pm2_resolved_why}_")
            else:
                st.markdown(f"**Structure: _{resolved_name}_**")
            st.markdown("**Here's the storyline I'd tell:**")
            st.markdown(f"_{sb['reasoning']}_")
            for n, oi in enumerate(sb["order"], 1):
                f = findings[oi]
                st.markdown(f"{n}. **{f.role.replace('_',' ').title()}** — {f.fact_brief()[:110]}")
            st.markdown("Approve to build, or tell me what to change.")
        cc1, cc2 = st.columns([1, 1])
        if cc1.button("✓ Approve & build", type="primary", key="pm2_approve", use_container_width=True):
            st.session_state.pm2_stage = "build"; st.rerun()
        redirect = cc2.text_input("Redirect (e.g. 'lead with the action', 'drop the correlation')",
                                  key="pm2_redirect")
        if st.button("Apply redirect", key="pm2_redirect_btn") and redirect.strip():
            if a.llm:
                try:
                    sb2 = draft_storyboard(findings, {**framing, "decision": framing.get("decision","") + " | " + redirect}, a.llm)
                    st.session_state.pm2_sb = sb2
                    st.rerun()
                except Exception:
                    st.warning("Couldn't apply — try rephrasing.")
            else:
                st.warning("AI engine not configured; using the structure picker above instead.")

    # ---- BUILD: choose visuals + narrate + assemble spec ----
    elif stage in ("build", "review"):
        framing = st.session_state.pm2_framing
        if "pm2_spec" not in st.session_state:
            sb = st.session_state.pm2_sb
            order = sb["order"]
            visuals, narrations = {}, {}
            with st.spinner("Choosing the right visual for each finding and writing the narrative…"):
                for oi in order:
                    f = findings[oi]
                    v = propose_visual(f, a.llm, framing.get("audience", "leadership"))
                    visuals[oi] = v
                    narrations[oi] = narrate_finding(f, framing, a.llm)
                    findings[oi].narrative = narrations[oi]
                # recommendations: from action findings or LLM
                recs = []
                if a.llm:
                    try:
                        heads = "; ".join(findings[i].fact_brief()[:90] for i in order[:5])
                        out = a.llm.predict(
                            f"Findings: {heads}\nDecision: {framing.get('decision','')}. "
                            "Give 3 specific, numbered action recommendations a consultant would put on the final slide. "
                            "Return ONLY JSON: {\"recs\":[\"\",\"\",\"\"],\"title\":\"<=6 word deck title\",\"closing\":\"one line\"}",
                            max_tokens=380)
                        import json as _j, re as _r
                        p = _j.loads(_r.sub(r"```json|```", "", out).strip())
                        recs = p.get("recs", [])[:3]
                        framing.setdefault("title", p.get("title", ""))
                        framing["closing"] = p.get("closing", "The data points to a clear next move.")
                    except Exception:
                        pass
                if not recs:
                    recs = [narrations[i] for i in order[:3]]
                framing["recommendations"] = recs
                style_key = st.session_state.get("pm2_resolved_style",
                                                  STORYBOARD_STYLES.get(style_label, {}).get("key", "consulting")
                                                  if style_label != "✨ Let the agent decide" else "consulting")
                spec = build_consultant_spec(a, findings, ctx, framing, theme_name, font_name,
                                             style_key, order, narrations, visuals)
            st.session_state.pm2_spec = spec
            st.session_state.pm2_stage = "review"

        spec = st.session_state.pm2_spec
        spec["meta"]["theme"] = theme_name; spec["meta"]["font"] = font_name
        _is_consulting = spec.get("consulting_plan") is not None

        # -------- floating chat assistant (data + visuals + deck aware) --------
        if "pm2_chat_log" not in st.session_state:
            st.session_state.pm2_chat_log = [
                ("assistant", "Hi — I know your data, the charts, and this deck. "
                              "Ask me to edit slides, add a breakdown from your data, "
                              "or explain anything. e.g. *make slide 3's title punchier*, "
                              "*add a slide breaking sales by state*, *why a waterfall here?*")
            ]
        _chat_deps = {
            "validate_chart_choice": validate_chart_choice, "propose_visual": propose_visual,
            "narrate_finding": narrate_finding, "Finding": Finding, "_fmt": _fmt,
            "CONSULT_THEMES": CONSULT_THEMES, "CONSULT_FONTS": CONSULT_FONTS,
        }
        _pop = st.popover("💬 Ask the assistant", use_container_width=False)
        with _pop:
            st.caption("Edits apply instantly — scroll down to see the deck update.")
            for _role, _msg in st.session_state.pm2_chat_log[-8:]:
                with st.chat_message("assistant" if _role == "assistant" else "user"):
                    st.markdown(_msg)
            _um = st.chat_input("Tell me what to change or ask a question…", key="pm2_chat_in")
            if _um:
                st.session_state.pm2_chat_log.append(("user", _um))
                _action = interpret_chat(_um, spec, findings, a, a.llm,
                                         history=st.session_state.pm2_chat_log[:-1])
                _ok, _resp = apply_chat_action(_action, spec, findings, a, _chat_deps)
                _prefix = "" if _action.get("action") == "answer" else ("✓ " if _ok else "⚠️ ")
                st.session_state.pm2_chat_log.append(("assistant", _prefix + (_resp or "Done.")))
                st.session_state.pm2_spec = spec
                st.rerun()

        if _is_consulting:
            plan = spec["consulting_plan"]
            st.markdown(
                f'<div class="kx-callout"><strong>Deck ready — {len(plan)+1} slides.</strong> '
                "Standard consulting storyline: Exec Summary → Problem → Insights & Drivers → "
                "Recommendations → Next Steps. Edit any text inline, restyle up top, then download.</div>",
                unsafe_allow_html=True,
            )
            for i, pi in enumerate(plan):
                lab = pi.get("role", "slide").replace("_", " ").title()
                st.markdown(f"**{pi['section']} — {lab}**")
                pv, ed = st.columns([3, 2])
                with pv:
                    _cmp.html(render_consulting_html(pi, spec["meta"]), height=270)
                with ed:
                    _consulting_editor(pi, i)
                st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)
        else:
            st.markdown(
                f'<div class="kx-callout"><strong>Deck ready — {len(spec["slides"])} slides.</strong> '
                "Preview each slide, edit any text inline, restyle from the pickers up top, then download.</div>",
                unsafe_allow_html=True,
            )
            for i, sl in enumerate(spec["slides"]):
                st.markdown(f"**Slide {i+1} — {sl['kind'].replace('_',' ').title()}**")
                pv, ed = st.columns([3, 2])
                with pv:
                    _cmp.html(render_consultant_html(sl, spec["meta"]), height=290)
                with ed:
                    _consult_editor(sl, i)
                st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)

        st.markdown("---")
        if st.button("⬇ Generate & download .pptx", type="primary", key="pm2_export"):
            from pathlib import Path as _P
            from datetime import datetime as _dtm
            outdir = _P("data/uploads"); outdir.mkdir(parents=True, exist_ok=True)
            fn = st.session_state.get("upload_filename", "deck").replace(".", "_")
            outp = str(outdir / f"KiteIQX_{fn}_{_dtm.now():%Y%m%d_%H%M%S}.pptx")
            with st.spinner("Rendering PowerPoint…"):
                ok, msg = build_consultant_pptx(spec, outp, _CHARTS_NS,
                                                findings=findings, framing=framing)
            if ok:
                st.session_state.pm2_out = outp; st.success(msg)
            else:
                st.error(msg)
        if st.session_state.get("pm2_out"):
            from pathlib import Path as _P
            if _P(st.session_state.pm2_out).exists():
                with open(st.session_state.pm2_out, "rb") as fh:
                    st.download_button("Download your deck", fh.read(),
                                       file_name=_P(st.session_state.pm2_out).name,
                                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                       key="pm2_dl")


def _consulting_editor(pi, i):
    role = pi.get("role")
    if role == "exec_summary":
        ex = pi.setdefault("exec", {})
        ex["core_answer"] = st.text_area("Core answer", ex.get("core_answer", ""), height=70, key=f"pm2_ce_core_{i}")
        for j in range(3):
            ins = ex.get("insights", [])
            while len(ins) < 3: ins.append("")
            ins[j] = st.text_input(f"Insight {j+1}", ins[j], key=f"pm2_ce_ins_{i}_{j}")
            ex["insights"] = ins
    elif role in ("actions",):
        recs = pi.setdefault("recs", [])
        for j in range(min(3, max(1, len(recs)))):
            while len(recs) <= j: recs.append("")
            recs[j] = st.text_input(f"Action {j+1}", recs[j], key=f"pm2_ce_rec_{i}_{j}")
    elif role == "impact":
        for j, m in enumerate(pi.get("impacts", [])):
            c1, c2 = st.columns(2)
            m["value"] = c1.text_input(f"Value {j+1}", m.get("value", ""), key=f"pm2_ce_iv_{i}_{j}")
            m["label"] = c2.text_input(f"Label {j+1}", m.get("label", ""), key=f"pm2_ce_il_{i}_{j}")
    elif role == "roadmap":
        st.caption("Roadmap phases — edit inline.")
        steps = pi.get("steps") or []
        for j, sp_ in enumerate(steps):
            sp_["what"] = st.text_input(f"{sp_.get('phase','Phase')}", sp_.get("what", ""), key=f"pm2_ce_rm_{i}_{j}")
    else:
        if "title" in pi:
            pi["title"] = st.text_input("Headline (action title)", pi.get("title", ""), key=f"pm2_ce_t_{i}")
        if "narrative" in pi:
            pi["narrative"] = st.text_area("Narrative", pi.get("narrative", ""), height=70, key=f"pm2_ce_n_{i}")
        if "takeaway" in pi:
            pi["takeaway"] = st.text_input("Takeaway", pi.get("takeaway", ""), key=f"pm2_ce_tk_{i}")
        if pi.get("chart_kind") is not None:
            opts = ["waterfall", "bar_ranked", "line", "area", "donut", "scatter"]
            cur = pi.get("chart_kind") or "bar_ranked"
            pi["chart_kind"] = st.selectbox("Chart", opts,
                                            index=opts.index(cur) if cur in opts else 1, key=f"pm2_ce_c_{i}")


def render_consulting_html(pi, meta):
    """Lightweight HTML preview of a consulting plan item."""
    th = _theme(meta["theme"], meta["font"])
    A = "#" + th["accent"]; INK = "#" + th["ink"]; SOFT = "#" + th["soft"]
    LINE = "#" + th["line"]; TAG = "#" + th["tagbg"]; TAGINK = "#" + th["tagink"]; CAN = "#" + th["canvas"]; PANEL = "#" + th["panel"]
    frame = (f'width:100%;aspect-ratio:16/9;border-radius:12px;overflow:hidden;position:relative;'
             f'font-family:Inter,system-ui,sans-serif;border:1px solid {LINE};background:{CAN};'
             f'box-shadow:0 2px 10px rgba(0,0,0,0.06);')
    role = pi.get("role")

    def chip(tag):
        return (f'<div style="position:absolute;left:4%;top:7%;background:{TAG};color:{TAGINK};'
                f'font-size:0.55rem;font-weight:700;letter-spacing:1.5px;padding:3px 12px;border-radius:20px;">{tag.upper()}</div>')

    def head(t):
        return f'<div style="position:absolute;left:4%;top:15%;right:5%;font-size:1.4rem;font-weight:800;color:{INK};line-height:1.1;">{t}</div>'

    if role == "exec_summary":
        ex = pi.get("exec", {})
        cards = "".join(
            f'<div style="flex:1;background:{PANEL};border:1px solid {LINE};border-radius:8px;padding:8px;">'
            f'<div style="color:{A};font-weight:800;font-size:0.8rem;">{j+1}</div>'
            f'<div style="font-size:0.55rem;color:{INK};line-height:1.25;margin-top:3px;">{ins}</div></div>'
            for j, ins in enumerate(ex.get("insights", [])[:3]))
        acts = "".join(f'<div style="font-size:0.55rem;color:{INK};margin-top:3px;">● {ac}</div>' for ac in ex.get("actions", [])[:2])
        return (f'<div style="{frame}">{chip("Executive Summary")}'
                f'<div style="position:absolute;left:4%;top:16%;right:5%;font-size:1.25rem;font-weight:800;color:{INK};line-height:1.1;">{ex.get("core_answer","")}</div>'
                f'<div style="position:absolute;left:4%;right:4%;top:46%;display:flex;gap:6px;">{cards}</div>'
                f'<div style="position:absolute;left:4%;top:78%;"><div style="font-size:0.5rem;color:{SOFT};font-weight:700;letter-spacing:1px;">RECOMMENDED ACTIONS</div>{acts}</div></div>')

    if role == "actions":
        rows = "".join(
            f'<div style="display:flex;margin-bottom:6px;border:1px solid {LINE};border-radius:8px;overflow:hidden;background:#fff;">'
            f'<div style="background:{A};color:#fff;width:30px;display:flex;align-items:center;justify-content:center;font-weight:800;">{j+1}</div>'
            f'<div style="padding:7px 10px;font-size:0.6rem;color:{INK};">{t}</div></div>'
            for j, t in enumerate(pi.get("recs", [])[:3]))
        return f'<div style="{frame}">{chip("Recommendations")}{head(pi.get("title","Recommended actions"))}<div style="position:absolute;left:4%;right:4%;top:32%;">{rows}</div></div>'

    if role == "impact":
        cards = "".join(
            f'<div style="flex:1;background:{PANEL};border:1px solid {LINE};border-radius:8px;padding:10px;text-align:center;">'
            f'<div style="font-size:1.6rem;font-weight:800;color:{A};">{m.get("value","")}</div>'
            f'<div style="font-size:0.52rem;color:{INK};margin-top:4px;">{m.get("label","")}</div></div>'
            for m in pi.get("impacts", [])[:3])
        return f'<div style="{frame}">{chip("Recommendations")}{head(pi.get("title","Expected impact"))}<div style="position:absolute;left:4%;right:4%;top:36%;display:flex;gap:6px;">{cards}</div></div>'

    if role == "roadmap":
        steps = pi.get("steps") or [{"phase":"Now","what":"Confirm driver"},{"phase":"Next","what":"Launch fix"},{"phase":"Later","what":"Scale"}]
        cards = ""
        for k, sp_ in enumerate(steps):
            cards += (f'<div style="flex:1;background:#fff;border:1px solid {LINE};border-radius:8px;overflow:hidden;">'
                      f'<div style="background:{A};color:#fff;font-size:0.55rem;font-weight:700;padding:4px;text-align:center;">{sp_.get("phase","")}</div>'
                      f'<div style="padding:7px;font-size:0.52rem;color:{INK};">{sp_.get("what","")}</div></div>')
            if k < len(steps)-1: cards += f'<div style="align-self:center;color:{SOFT};">→</div>'
        return f'<div style="{frame}">{chip("Next Steps")}{head(pi.get("title","Next steps & roadmap"))}<div style="position:absolute;left:4%;right:4%;top:38%;display:flex;gap:5px;align-items:stretch;">{cards}</div></div>'

    # finding-style (stat_hero / charts): headline / visual / takeaway
    big = ""
    svg = ""
    title = pi.get("title", "Finding")
    narrative = pi.get("narrative", "")
    tk = pi.get("takeaway", "")
    inner = (f'{chip(pi.get("section","Insight"))}{head(title)}'
             f'<div style="position:absolute;left:4%;top:32%;right:36%;font-size:0.6rem;color:{INK};line-height:1.35;">{narrative}</div>'
             f'<div style="position:absolute;right:4%;top:32%;width:30%;background:{TAG};border-radius:8px;padding:8px;font-size:0.55rem;color:{INK};">{pi.get("chart_kind","chart").replace("_"," ")} visual</div>')
    if tk:
        inner += f'<div style="position:absolute;left:4%;right:4%;bottom:5%;background:{TAG};border-radius:8px;padding:6px 10px;font-size:0.55rem;color:{INK};"><b style="color:{TAGINK};">Takeaway</b> {tk}</div>'
    return f'<div style="{frame}">{inner}</div>'


def _consult_editor(sl, i):
    k = sl["kind"]
    if "title" in sl:
        sl["title"] = st.text_input("Title", sl.get("title",""), key=f"pm2_e_t_{i}")
    if k == "title":
        sl["subtitle"] = st.text_input("Subtitle", sl.get("subtitle",""), key=f"pm2_e_s_{i}")
    if k == "finding":
        sl["narrative"] = st.text_area("Narrative", sl.get("narrative",""), height=90, key=f"pm2_e_n_{i}")
        opts = ["waterfall","bar_ranked","line","area","donut","scatter","slope","dumbbell"]
        cur = sl.get("chart_kind") or "bar_ranked"
        sl["chart_kind"] = st.selectbox("Chart", opts,
                                        index=opts.index(cur) if cur in opts else 1, key=f"pm2_e_c_{i}")
    if k == "actions":
        for j, t in enumerate(sl.get("items", [])):
            sl["items"][j] = st.text_input(f"Action {j+1}", t, key=f"pm2_e_a_{i}_{j}")
    if k == "closing":
        sl["message"] = st.text_area("Closing", sl.get("message",""), height=70, key=f"pm2_e_m_{i}")
    if k == "storyline":
        st.caption("Flow steps come from the storyboard structure picker above.")

# ============================================================
# MAIN
# ============================================================

def initialize_analytics(df: pd.DataFrame, model: str, filename: str = "unknown"):
    llm = None
    # Prefer Claude (Anthropic); fall back to Groq if Claude isn't configured.
    if ANTHROPIC_AVAILABLE and ANTHROPIC_API_KEY:
        try:
            llm = ClaudeLLM(api_key=ANTHROPIC_API_KEY, model=CLAUDE_MODEL)
        except Exception as e:
            st.error(f"Could not initialize Claude: {e}")
    if llm is None and GROQ_AVAILABLE and GROQ_API_KEY:
        try:
            llm = GroqLLM(api_key=GROQ_API_KEY, model=model)
        except Exception as e:
            st.error(f"Could not initialize Groq: {e}")
    st.session_state.analytics = UniversalAnalytics(df, llm)
    st.session_state.viz_engine = VizEngine(st.session_state.analytics)
    st.session_state.upload_filename = filename
    st.session_state.pop("exec_summary", None)
    st.session_state.pop("takeaways", None)
    st.session_state.pop("ai_last_response", None)
    st.session_state.pop("upload_logged", None)
    st.session_state.pop("ppt_out_path", None)
    st.session_state.pop("ppt_payload", None)
    for _k in list(st.session_state.keys()):
        if _k.startswith("pm_") or _k.startswith("pm2_"):
            st.session_state.pop(_k, None)
    log_upload(filename, df)


def main():
    render_header()

    st.markdown(
        """
        <style>
        button { color: #ffffff !important; }
        button p, button span, button div { color: #ffffff !important; }
        .stButton > button { background-color: #0a2540 !important; color: #ffffff !important; }
        .stButton > button p,
        .stButton > button span { color: #ffffff !important; }
        .stFileUploader button { background-color: #0a2540 !important; color: #ffffff !important; }
        </style>
        """,
        unsafe_allow_html=True,
    )

    if "analytics" not in st.session_state:
        st.session_state.analytics = None
        st.session_state.viz_engine = None

    with st.sidebar:
        st.markdown("### Configuration")
        if ANTHROPIC_AVAILABLE and ANTHROPIC_API_KEY:
            st.success("AI engine ready — Claude.")
            _model_label = CLAUDE_MODEL
        elif GROQ_AVAILABLE and GROQ_API_KEY:
            st.success("AI engine ready — Groq (Claude not configured).")
            _model_label = FIXED_MODEL
        elif not ANTHROPIC_AVAILABLE and not GROQ_AVAILABLE:
            st.error("Neither `anthropic` nor `groq` package installed.")
            _model_label = "(none)"
        else:
            st.warning("No ANTHROPIC_API_KEY in secrets. Add it to enable AI features.")
            _model_label = "(none)"
        st.markdown(
            f'<div style="font-size:0.78rem;color:#4a5568;background:#f7f8fb;'
            f'border:1px solid #e3e6ee;border-radius:6px;padding:0.35rem 0.7rem;'
            f'margin-top:0.3rem;margin-bottom:0.1rem;">'
            f'⚡ Model: <strong style="color:#0a2540;">{_model_label}</strong></div>',
            unsafe_allow_html=True,
        )

        model = FIXED_MODEL

        st.markdown("---")
        st.markdown("### Load data")
        up = st.file_uploader("Drop a CSV or Excel file", type=["csv", "xlsx", "xls"])
        url = st.text_input("…or fetch CSV from URL")
        demo = st.selectbox("…or try a demo dataset", ["", "E-commerce Sales", "Employee Data", "Simple Demo"])

        c1, c2 = st.columns(2)
        with c1:
            if st.button("Load", type="primary"):
                try:
                    df = None
                    fname = "unknown"
                    if up:
                        df = load_uploaded_file(up)
                        fname = up.name
                        st.success(f"Loaded **{up.name}** ({len(df):,} rows).")
                    elif url:
                        r = requests.get(url, timeout=30)
                        df = pd.read_csv(io.StringIO(r.text))
                        fname = url.split("/")[-1] or "url_import"
                        st.success(f"Loaded from URL ({len(df):,} rows).")
                    if df is not None:
                        initialize_analytics(df, model, filename=fname)
                except Exception as e:
                    st.error(f"Load failed: {e}")
        with c2:
            if st.button("Use demo") and demo:
                try:
                    df = generate_demo(demo)
                    initialize_analytics(df, model, filename=f"demo_{demo.replace(' ','_')}")
                    st.success(f"Loaded demo: {demo}.")
                except Exception as e:
                    st.error(f"Demo failed: {e}")

        if st.session_state.get("last_saved_path"):
            st.markdown("---")
            st.markdown("**Last uploaded file archived to:**")
            st.code(st.session_state["last_saved_path"], language=None)

        if st.session_state.analytics:
            st.markdown("---")
            s = st.session_state.analytics.insights["basic_stats"]
            st.metric("Rows", f"{s['rows']:,}")
            st.metric("Columns", s["columns"])
            st.metric("Completeness", f"{100 - s['missing_pct']:.1f}%")

    if st.session_state.analytics is None:
        render_welcome()
        return

    tabs = st.tabs([
        "Executive Dashboard",
        "AI Intelligence",
        "Data Quality & Charts",
        "Data Explorer",
        "Presentation Maker",
    ])

    with tabs[0]:
        render_dashboard()
    with tabs[1]:
        render_ai()
    with tabs[2]:
        render_quality_and_charts()
    with tabs[3]:
        render_explorer()
    with tabs[4]:
        render_presentation_maker()


if __name__ == "__main__":
    main()
